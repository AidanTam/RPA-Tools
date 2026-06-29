# ##### Streamlit app – Z‑Score QAQC Charts (full version, CRM‑legend fixed)
# """
# Author: Hamza Salhi – July 2025
# Interactive QAQC tool that plots Z‑scores for CRMs and highlights ±2 σ / ±3 σ limits.

# Key features
# ------------
# • Upload CSV → map columns → filter Elements / CRMs / Companies / Dates.
# • Z‑score scatter vs sequential order (colour = CRM).
# • Control lines at ±3 σ (green), ±2 σ (lime), 0 (black, dashed).
# • Legend lists *all* CRMs (no limit) plus the control‑line labels.
# • Sidebar style controls, summary table, Excel & PPT exports, JSON config save.
# """

# # ───────── imports ─────────
# import os, io, json, tempfile, re, random, datetime as dt
# import numpy as np, pandas as pd
# import matplotlib.pyplot as plt, seaborn as sns
# from matplotlib.lines import Line2D
# from matplotlib.ticker import FuncFormatter
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches

# # ───────── page ─────────
# st.set_page_config(page_title="Z‑Score QAQC Charts", layout="wide")
# st.title("🧮 Z‑Score QAQC Charts")

# # ───────── helper utils ─────────
# _safe_pat = re.compile(r"[^0-9A-Za-z_]+")
# idx_of    = lambda v, opts: list(opts).index(v) if v in opts else 0
# _fmt      = FuncFormatter(lambda v, _: f"{v:.1f}")

# def safe_defaults(defs, opts): return [d for d in (defs or []) if d in opts]

# def clean_numeric(s):
#     if pd.api.types.is_numeric_dtype(s):
#         return pd.to_numeric(s, errors="coerce")
#     return pd.to_numeric(
#         s.astype(str)
#          .str.strip()
#          .str.replace(r"[^\d\.\-eE+]", "", regex=True),
#         errors="coerce"
#     )

# def tmp_fig(fig):
#     t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(t.name, dpi=300, bbox_inches="tight"); plt.close(fig); return t.name

# def shuffled_palette(base="tab10", n=10, seed=0):
#     """Return *n* colours from *base* palette, shuffled with *seed*."""
#     pal = sns.color_palette(base, n)
#     if seed:
#         rng = random.Random(seed); rng.shuffle(pal)
#     return pal

# # ───────── config load ─────────
# cfg_up = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
# if cfg_up and not st.session_state.get("cfg_loaded"):
#     st.session_state.update(json.load(cfg_up)); st.session_state.cfg_loaded=True
#     st.success("Config loaded!")

# # ───────── CSV load ─────────
# csv_up = st.sidebar.file_uploader("📂 Upload CSV", ["csv"], key="csv")
# df=None
# if csv_up is not None:
#     df = pd.read_csv(csv_up)
#     st.session_state["csv_path"] = ""
# elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
#     df = pd.read_csv(st.session_state["csv_path"])
# if df is None:
#     st.info("Upload a CSV or load a config to begin."); st.stop()

# # ───────── column mapping ─────────
# with st.sidebar.expander("🗺️ Column mapping", True):
#     elem_col = st.selectbox("Element column", df.columns,
#                              index=idx_of(st.session_state.get("elem_col",df.columns[0]), df.columns))
#     crm_col  = st.selectbox("CRM column",    df.columns,
#                              index=idx_of(st.session_state.get("crm_col",df.columns[1]), df.columns))
#     val_col  = st.selectbox("Value column",  df.columns,
#                              index=idx_of(st.session_state.get("val_col",df.columns[2]), df.columns))
#     exp_col  = st.selectbox("Expected value column", df.columns,
#                              index=idx_of(st.session_state.get("exp_col",df.columns[3]), df.columns))
#     sd_col   = st.selectbox("SD column", df.columns,
#                              index=idx_of(st.session_state.get("sd_col",df.columns[4]), df.columns))
#     date_col = st.selectbox("Date column", df.columns,
#                              index=idx_of(st.session_state.get("date_col",df.columns[5]), df.columns))
#     comp_col = st.selectbox("Company column", df.columns,
#                              index=idx_of(st.session_state.get("comp_col",df.columns[6]), df.columns))

# # ───────── filters ─────────
# with st.sidebar.expander("🔎 Core filters"):
#     elem_opts = df[elem_col].dropna().unique()
#     elems_sel = st.multiselect("Elements", elem_opts,
#                                default=safe_defaults(st.session_state.get("elems_sel"), elem_opts),
#                                key="elems_sel")

#     crm_opts = df[crm_col].dropna().unique()
#     crms_sel = st.multiselect("CRMs", crm_opts,
#                               default=safe_defaults(st.session_state.get("crms_sel"), crm_opts),
#                               key="crms_sel")

#     comp_opts = df[comp_col].dropna().unique()
#     comps_sel = st.multiselect("Companies", comp_opts,
#                                default=safe_defaults(st.session_state.get("comps_sel"), comp_opts),
#                                key="comps_sel")

# with st.sidebar.expander("🗓️ Date filter"):
#     df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
#     dmin, dmax = df[date_col].min(), df[date_col].max()
#     date_rng = st.date_input("Range",
#                              st.session_state.get("date_rng",
#                                      (None if pd.isna(dmin) else dmin.date(),
#                                       None if pd.isna(dmax) else dmax.date())),
#                              key="date_rng")

# # ───────── style ─────────
# with st.sidebar.expander("📐 Style"):
#     point_size  = st.number_input("Point size (pts)", 1,
#                                   value=int(st.session_state.get("p_size", 8)), step=1,
#                                   key="p_size")
#     line_width  = st.number_input("Line width", 0.1,
#                                   value=float(st.session_state.get("l_width", 0.8)), step=0.1,
#                                   format="%.2f", key="l_width")
#     palette_seed = st.number_input("Palette seed", 0, 1000,
#                                    value=int(st.session_state.get("pal_seed", 0)),
#                                    step=1, key="pal_seed")

# # ───────── run trigger ─────────
# if st.sidebar.button("🚀 Run..."): st.session_state.run_plots = True
# if not st.session_state.get("run_plots"):
#     st.info("Click **Run** to generate charts."); st.stop()

# # ───────── preprocess ─────────
# for c in (val_col, exp_col, sd_col):
#     df[c] = clean_numeric(df[c])
# df = df.dropna(subset=[val_col, exp_col, sd_col, date_col])
# df["z_score"] = (df[val_col] - df[exp_col]) / df[sd_col]
# df = df.sort_values(date_col).reset_index(drop=True)
# df["SortIndex"] = df.index

# # ───────── crunch & plot ─────────
# ppt, summary_rows = Presentation(), []
# size_pts2 = point_size ** 2   # scatter expects area
# ctrl_lines = [(3, "+3 SD", "#7FA37F"), (-3, "-3 SD", "#7FA37F"),
#               (2, "+2 SD", "#B5ED38"), (-2, "-2 SD", "#B5ED38"),
#               (0, "0", "#000000")]

# for el in elems_sel:
#     dfe = df[df[elem_col] == el]
#     if crms_sel:
#         dfe = dfe[dfe[crm_col].isin(crms_sel)]
#     if comps_sel:
#         dfe = dfe[dfe[comp_col].isin(comps_sel)]
#     if date_rng:
#         start, end = map(pd.to_datetime, date_rng); end += pd.Timedelta(days=1)
#         dfe = dfe[(dfe[date_col] >= start) & (dfe[date_col] < end)]
#     if dfe.empty:
#         continue

#     # palette that covers every CRM
#     crm_list = sorted(dfe[crm_col].unique())
#     base_pal = shuffled_palette("tab10", 10, palette_seed)
#     colors = (base_pal * ((len(crm_list) // len(base_pal)) + 1))[:len(crm_list)]
#     crm_palette = dict(zip(crm_list, colors))

#     fig, ax = plt.subplots(figsize=(15, 6))

#     # control lines
#     for y, lbl, col in ctrl_lines:
#         ax.axhline(y=y, color=col, lw=line_width, ls="-" if y else "--")

#     # scatterplot
#     sns.scatterplot(
#         data=dfe,
#         x="SortIndex",
#         y="z_score",
#         hue=crm_col,
#         palette=crm_palette,
#         s=size_pts2,
#         edgecolor=None,
#         ax=ax,
#         legend=False,
#     )

#     # x‑tick labels
#     tick_spacing = max(1, len(dfe) // 20)
#     ax.set_xticks(dfe["SortIndex"][::tick_spacing])
#     ax.set_xticklabels(
#         dfe[date_col].dt.strftime("%Y-%m")[::tick_spacing],
#         rotation=45,
#         ha="right",
#     )

#     # company bands
#     for comp in dfe[comp_col].unique():
#         sub = dfe[dfe[comp_col] == comp]
#         mi, ma = sub["SortIndex"].min(), sub["SortIndex"].max()
#         ax.axvline(mi, ls="--", lw=0.8, color="lightgrey")
#         ax.axvline(ma, ls="--", lw=0.8, color="lightgrey")
#         ax.annotate(
#             comp,
#             xy=((mi + ma) / 2, 26),
#             ha="center",
#             bbox=dict(fc="white", ec="black", pad=0.3),
#             fontsize=9,
#         )

#     ax.set_ylim(-30, 30)
#     ax.set_ylabel("Z‑score")
#     ax.set_xlabel("Date")
#     ax.set_title(f"Z‑Score by CRM & Company – {el}")

#     # legend: all CRMs + control lines
#     crm_handles = [
#         Line2D(
#             [],
#             [],
#             marker="o",
#             color="none",
#             markerfacecolor=crm_palette[c],
#             markersize=6,
#             label=c,
#         )
#         for c in crm_list
#     ]
#     line_handles = [
#         Line2D([0, 1], [0, 0], color=col, lw=1.5, ls="-" if y else "--", label=lbl)
#         for y, lbl, col in ctrl_lines
#     ]
#     ax.legend(
#         crm_handles + line_handles,
#         [h.get_label() for h in crm_handles + line_handles],
#         title="CRM / Limits",
#         bbox_to_anchor=(0.5, -0.28),
#         loc="upper center",
#         ncol=7,
#         frameon=True,
#         fontsize=8,
#         handlelength=3,
#     )

#     plt.tight_layout()
#     plt.subplots_adjust(bottom=0.4)  # space for legend

#     # Streamlit + PPT
#     st.pyplot(fig)
#     slide = ppt.slides.add_slide(ppt.slide_layouts[5])
#     slide.shapes.add_picture(tmp_fig(fig), Inches(0.5), Inches(0.8), width=Inches(9))

#     # summary stats
#     n = len(dfe)
#     out2 = (np.abs(dfe["z_score"]) > 2).sum()
#     out3 = (np.abs(dfe["z_score"]) > 3).sum()
#     summary_rows.append(
#         {
#             "Element": el,
#             "Samples": n,
#             "|z|>2": out2,
#             "Rate_>2_%": round(out2 / n * 100, 1),
#             "|z|>3": out3,
#             "Rate_>3_%": round(out3 / n * 100, 1),
#         }
#     )

# # ───────── summary & exports ─────────
# summary_df = pd.DataFrame(summary_rows)
# st.subheader("📊 Summary")
# st.dataframe(summary_df)

# xls = io.BytesIO()
# summary_df.to_excel(xls, index=False)
# st.download_button("Download summary.xlsx", xls.getvalue(), file_name="zscore_summary.xlsx")

# ppt_bytes = io.BytesIO()
# ppt.save(ppt_bytes)
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="zscore_plots.pptx")

# # ───────── save config ─────────
# serial = {}
# for k, v in st.session_state.items():
#     if k.startswith("_"):
#         continue
#     try:
#         json.dumps(v)
#         serial[k] = v
#     except TypeError:
#         pass
# cfg_bytes = io.BytesIO(json.dumps(serial, indent=2).encode())
# st.sidebar.download_button(
#     "💾 Download config JSON",
#     cfg_bytes.getvalue(),
#     file_name="zscore_config.json",
#     mime="application/json",
# )


### Version 2 #####
##### Streamlit app – Z‑Score QAQC Charts (with Y‑axis zoom)
"""
Author: Hamza Salhi – July 2025
Z‑score QAQC tool with column mapping, filters, full CRM legend, and
user‑defined Y‑axis (Z‑score) limits.

How to use
----------
1. `streamlit run zscore_qaqc_app.py`
2. Upload CSV → map columns.
3. In *Style* panel set *Min Z* / *Max Z* to zoom (default –30 to 30).
4. Click **Run** to generate charts.
"""

# ───────── imports ─────────
import os, io, json, tempfile, re, random, datetime as dt
import numpy as np, pandas as pd
import matplotlib.pyplot as plt, seaborn as sns
from matplotlib.lines import Line2D
from matplotlib.ticker import FuncFormatter
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pyrpa.UI import common

# ───────── page ─────────
st.set_page_config(page_title="Z‑Score QAQC Charts", layout="wide")
st.title("🧮 Z‑Score QAQC Charts")

# ───────── helpers ─────────
_safe = re.compile(r"[^0-9A-Za-z_]+")
idx_of  = lambda v, opts: list(opts).index(v) if v in opts else 0
_fmt    = FuncFormatter(lambda v, _: f"{v:.1f}")

def safe_defaults(defs, opts): return [d for d in (defs or []) if d in opts]

def clean_numeric(s):
    if pd.api.types.is_numeric_dtype(s):
        return pd.to_numeric(s, errors="coerce")
    return pd.to_numeric(
        s.astype(str).str.replace(r"[^\d\.\-eE+]", "", regex=True),
        errors="coerce",
    )

def tmp_fig(fig):
    f = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(f.name, dpi=300, bbox_inches="tight"); plt.close(fig); return f.name

def shuffled_palette(base="tab10", n=10, seed=0):
    pal = sns.color_palette(base, n)
    if seed:
        rng = random.Random(seed); rng.shuffle(pal)
    return pal

# ───────── config load ─────────
cfg_up = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
if cfg_up and not st.session_state.get("cfg_loaded"):
    st.session_state.update(json.load(cfg_up)); st.session_state.cfg_loaded=True
    st.success("Config loaded!")

# ───────── data load ─────────
csv_up = st.sidebar.file_uploader("📂 Upload data (CSV or Excel)", ["csv","xlsx","xls"], key="csv")
df = None
if csv_up is not None:
    df = common.read_data_file(csv_up); st.session_state["csv_path"] = ""
elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
    df = common.read_data_file(st.session_state["csv_path"])
if df is None:
    st.info("Upload CSV or load config to begin."); st.stop()

# ───────── column mapping ─────────
with st.sidebar.expander("🗺️ Column mapping", True):
    elem_col = st.selectbox("Element column", df.columns,
                             index=idx_of(st.session_state.get("elem_col", df.columns[0]), df.columns),
                             key="elem_col")
    crm_col  = st.selectbox("CRM column", df.columns,
                             index=idx_of(st.session_state.get("crm_col", df.columns[1]), df.columns),
                             key="crm_col")
    val_col  = st.selectbox("Value column", df.columns,
                             index=idx_of(st.session_state.get("val_col", df.columns[2]), df.columns),
                             key="val_col")
    exp_col  = st.selectbox("Expected value column", df.columns,
                             index=idx_of(st.session_state.get("exp_col", df.columns[3]), df.columns),
                             key="exp_col")
    sd_col   = st.selectbox("SD column", df.columns,
                             index=idx_of(st.session_state.get("sd_col", df.columns[4]), df.columns),
                             key="sd_col")
    date_col = st.selectbox("Date column", df.columns,
                             index=idx_of(st.session_state.get("date_col", df.columns[5]), df.columns),
                             key="date_col")
    comp_col = st.selectbox("Company column", df.columns,
                             index=idx_of(st.session_state.get("comp_col", df.columns[6]), df.columns),
                             key="comp_col")

# ───────── filters ─────────
with st.sidebar.expander("🔎 Core filters"):
    elem_opts = df[elem_col].dropna().unique()
    elems_sel = st.multiselect("Elements", elem_opts,
                               default=safe_defaults(st.session_state.get("elems_sel"), elem_opts),
                               key="elems_sel")

    crm_opts  = df[crm_col].dropna().unique()
    crms_sel  = st.multiselect("CRMs", crm_opts,
                               default=safe_defaults(st.session_state.get("crms_sel"), crm_opts),
                               key="crms_sel")

    comp_opts = df[comp_col].dropna().unique()
    comps_sel = st.multiselect("Companies", comp_opts,
                               default=safe_defaults(st.session_state.get("comps_sel"), comp_opts),
                               key="comps_sel")

with st.sidebar.expander("🗓️ Date filter"):
    df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
    dmin, dmax = df[date_col].min(), df[date_col].max()
    date_rng = st.date_input("Range",
                             st.session_state.get("date_rng",
                                     (None if pd.isna(dmin) else dmin.date(),
                                      None if pd.isna(dmax) else dmax.date())),
                             key="date_rng")

# ───────── style + Y‑axis limits ─────────
with st.sidebar.expander("📐 Style"):
    point_size  = st.number_input("Point size (pts)", 1,
                                  value=int(st.session_state.get("p_size", 8)), step=1,
                                  key="p_size")
    line_width  = st.number_input("Line width", 0.1,
                                  value=float(st.session_state.get("l_width", 0.8)), step=0.1,
                                  format="%.2f", key="l_width")
    palette_seed = st.number_input("Palette seed", 0, 1000,
                                   value=int(st.session_state.get("pal_seed", 0)), step=1,
                                   key="pal_seed")
    st.markdown("**Z‑Score limits (Y‑axis zoom)**")
    zmin = st.number_input("Min Z", -100.0, 0.0,
                           value=float(st.session_state.get("zmin", -30.0)),
                           step=1.0, key="zmin")
    zmax = st.number_input("Max Z", 0.0, 100.0,
                           value=float(st.session_state.get("zmax", 30.0)),
                           step=1.0, key="zmax")
    if zmin >= zmax:
        st.warning("⚠️ Min Z must be **less** than Max Z")

# ───────── run trigger ─────────
if st.sidebar.button("🚀 Run"): st.session_state.run_plots = True
if not st.session_state.get("run_plots"):
    st.info("Click **Run** to generate charts."); st.stop()

# ───────── preprocess ─────────
for c in (val_col, exp_col, sd_col):
    df[c] = clean_numeric(df[c])
df = df.dropna(subset=[val_col, exp_col, sd_col, date_col])
df["z_score"] = (df[val_col] - df[exp_col]) / df[sd_col]
df = df.sort_values(date_col).reset_index(drop=True)
df["SortIndex"] = df.index

# ───────── constants ─────────
ctrl_lines = [(3, "+3 SD", "#7FA37F"),
              (-3, "‑3 SD", "#7FA37F"),
              (2, "+2 SD", "#B5ED38"),
              (-2, "‑2 SD", "#B5ED38"),
              (0, "0", "#000000")]

ppt, summary_rows = Presentation(), []
size_pts2 = point_size ** 2

# ───────── crunch & plot ─────────
for el in elems_sel:
    sub = df[df[elem_col] == el]
    if crms_sel:
        sub = sub[sub[crm_col].isin(crms_sel)]
    if comps_sel:
        sub = sub[sub[comp_col].isin(comps_sel)]
    if date_rng:
        start, end = map(pd.to_datetime, date_rng); end += pd.Timedelta(days=1)
        sub = sub[(sub[date_col] >= start) & (sub[date_col] < end)]
    if sub.empty: continue

    # palette to cover all CRMs
    crms = sorted(sub[crm_col].unique())
    base_pal = shuffled_palette("tab10", 10, palette_seed)
    colors = (base_pal * ((len(crms) // len(base_pal)) + 1))[:len(crms)]
    crm_palette = dict(zip(crms, colors))

    fig, ax = plt.subplots(figsize=(15, 6))

    # control lines
    for y, lbl, col in ctrl_lines:
        ax.axhline(y=y, color=col, lw=line_width, ls="-" if y else "--")

    # scatter
    sns.scatterplot(
        data=sub,
        x="SortIndex", y="z_score",
        hue=crm_col, palette=crm_palette,
        s=size_pts2, edgecolor=None, ax=ax, legend=False,
    )

    # x‑ticks
    tick_spacing = max(1, len(sub) // 20)
    ax.set_xticks(sub["SortIndex"][::tick_spacing])
    ax.set_xticklabels(sub[date_col].dt.strftime("%Y-%m")[::tick_spacing],
                       rotation=45, ha="right")

    # company bands
    for comp in sub[comp_col].unique():
        cdf = sub[sub[comp_col] == comp]
        mi, ma = cdf["SortIndex"].min(), cdf["SortIndex"].max()
        ax.axvline(mi, ls="--", lw=0.8, color="lightgrey")
        ax.axvline(ma, ls="--", lw=0.8, color="lightgrey")
        ax.annotate(comp, xy=((mi + ma) / 2, zmax - (zmax - zmin) * 0.05),
                    ha="center", fontsize=9,
                    bbox=dict(fc="white", ec="black", pad=0.3))

    # axes & title
    ax.set_ylim(zmin, zmax)
    ax.set_ylabel("Z‑score")
    ax.set_xlabel("Date")
    ax.set_title(f"Z‑Score by CRM & Company – {el}")

    # legend (all CRMs + control lines)
    crm_handles = [Line2D([], [], marker="o", color="none",
                          markerfacecolor=crm_palette[c], markersize=6, label=c)
                   for c in crms]
    line_handles = [Line2D([0, 1], [0, 0], color=col, lw=1.5,
                           ls="-" if y else "--", label=lbl)
                    for y, lbl, col in ctrl_lines]
    ax.legend(crm_handles + line_handles,
              [h.get_label() for h in crm_handles + line_handles],
              title="CRM / Limits",
              bbox_to_anchor=(0.5, -0.28), loc="upper center",
              ncol=7, frameon=True, fontsize=8, handlelength=3)

    plt.tight_layout()
    plt.subplots_adjust(bottom=0.4)

    # Streamlit + PPT
    st.pyplot(fig)
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.add_picture(tmp_fig(fig), Inches(0.5), Inches(0.8), width=Inches(9))

    # summary row
    n = len(sub)
    out2 = (np.abs(sub["z_score"]) > 2).sum()
    out3 = (np.abs(sub["z_score"]) > 3).sum()
    summary_rows.append({
        "Element": el, "Samples": n,
        "|z|>2": out2, "Rate_>2_%": round(out2 / n * 100, 1),
        "|z|>3": out3, "Rate_>3_%": round(out3 / n * 100, 1),
    })

# ───────── summary & exports ─────────
summary_df = pd.DataFrame(summary_rows)
st.subheader("📊 Summary"); st.dataframe(summary_df)

xls = io.BytesIO(); summary_df.to_excel(xls, index=False)
st.download_button("Download summary.xlsx", xls.getvalue(), file_name="zscore_summary.xlsx")

ppt_bytes = io.BytesIO(); ppt.save(ppt_bytes)
st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="zscore_plots.pptx")

# ───────── save config ─────────
serial = {}
for k, v in st.session_state.items():
    if k.startswith("_"): continue
    try:
        json.dumps(v); serial[k] = v
    except TypeError:
        pass
cfg_bytes = io.BytesIO(json.dumps(serial, indent=2).encode())
st.sidebar.download_button("💾 Download config JSON", cfg_bytes.getvalue(),
                           file_name="zscore_config.json",
                           mime="application/json")
