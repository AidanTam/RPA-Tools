
# import os, io, tempfile, datetime as dt
# import numpy as np, pandas as pd, matplotlib.pyplot as plt, streamlit as st
# from pptx import Presentation; from pptx.util import Inches

# st.set_page_config(page_title="Duplicate QAQC Charts", layout="wide")
# st.title("🪓 Duplicate QAQC Charts")

# # ── 1. Upload CSV ───────────────────────────────────────────────────
# csv_file = st.sidebar.file_uploader("Upload duplicate CSV", ["csv"])
# if not csv_file:
#     st.stop()
# df = pd.read_csv(csv_file)

# def guess_col(name):
#     key = name.lower().replace(" ", "")
#     return next((c for c in df.columns if c.lower().replace(" ", "") == key), None)

# # ── 2. Column mapping ───────────────────────────────────────────────
# with st.sidebar.expander("🗺️  Column mapping", expanded=True):
#     lab_col = st.selectbox("Lab column", df.columns,
#         index=int(df.columns.get_indexer([guess_col("Lab")])[0]) if guess_col("Lab") else 0)
#     dup_col = st.selectbox("Duplicate-type column", df.columns,
#         index=int(df.columns.get_indexer([guess_col("Duplicate Type")])[0]) if guess_col("Duplicate Type") else 0)

# # ── 3. Core filters ─────────────────────────────────────────────────
# with st.sidebar.expander("🔎 Filter core columns"):
#     labs_sel = st.multiselect("Labs", df[lab_col].dropna().unique(), default=list(df[lab_col].dropna().unique()))
#     dup_sel  = st.multiselect("Duplicate types", df[dup_col].dropna().unique(), default=list(df[dup_col].dropna().unique()))

# # ── 4. Free-form categorical filters (x3) ───────────────────────────
# cat_filters = []
# with st.sidebar.expander("🎛️  Extra categorical filters"):
#     for i in range(1, 4):
#         col = st.selectbox(f"Filter {i} – column", ["None"]+list(df.columns), key=f"cat_col_{i}")
#         if col != "None":
#             vals = st.multiselect("Values", df[col].dropna().unique(), key=f"cat_vals_{i}")
#             cat_filters.append((col, vals))

# # ── 5. Date filter ──────────────────────────────────────────────────
# with st.sidebar.expander("🗓️  Date filter"):
#     date_col = st.selectbox("Date column", ["None"]+list(df.columns), index=0)
#     date_range = None
#     if date_col != "None":
#         df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
#         dmin, dmax = df[date_col].min().date(), df[date_col].max().date()
#         date_range = st.date_input("Date range", (dmin, dmax))

# # ── 6. Plot options (axis + log) ────────────────────────────────────
# with st.sidebar.expander("📐 Plot options"):
#     use_log = st.checkbox("Log scale (scatter axes)")
#     custom_lim = st.checkbox("Custom X/Y limits")
#     x_min = y_min = x_max = y_max = None
#     if custom_lim:
#         x_min = st.number_input("X min", value=0.0)
#         x_max = st.number_input("X max", value=0.0)
#         y_min = st.number_input("Y min", value=0.0)
#         y_max = st.number_input("Y max", value=0.0)

# # ── 7. Element-pair detection ───────────────────────────────────────
# elements_info = {}
# for c in df.columns:
#     if c.endswith("_or"):
#         base = c[:-3]; dupc = f"{base}_dp"
#         if dupc in df.columns:
#             unit = "ppm" if "ppm" in base else "%" if "%" in base else ""
#             el   = base.split("_")[0].replace("(ppm)","").replace("(%)","")
#             elements_info[el] = {"unit":unit, "original_col":c, "duplicate_col":dupc}

# elements_sel = st.sidebar.multiselect("Elements", list(elements_info), default=list(elements_info))
# run = st.sidebar.button("🚀 Run analysis")
# if not run: st.stop()

# # ── 8. Prep outputs ────────────────────────────────────────────────
# ppt, summary_rows = Presentation(), []
# hard_pct = {"PD":.1, "CD":.2, "FD":.3}
# def tmp_fig(fig):
#     t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(t.name, dpi=300, bbox_inches="tight"); plt.close(fig); return t.name

# # ── 9. Main crunch loop ─────────────────────────────────────────────
# for el in elements_sel:
#     info = elements_info[el]
#     base = df[(df[lab_col].isin(labs_sel)) & (df[dup_col].isin(dup_sel))]
#     for col, vals in cat_filters:
#         if vals: base = base[base[col].isin(vals)]
#     if date_col != "None" and date_range:
#         start, end = pd.to_datetime(date_range[0]), pd.to_datetime(date_range[1]) + pd.Timedelta(days=1)
#         base = base[(base[date_col] >= start) & (base[date_col] < end)]

#     for lab in base[lab_col].unique():
#         for dtyp in base[dup_col].unique():
#             sub = base[(base[lab_col]==lab) & (base[dup_col]==dtyp)].dropna(subset=[info["original_col"], info["duplicate_col"]])
#             if sub.empty: continue

#             orig, dup = sub[info["original_col"]], sub[info["duplicate_col"]]
#             hard = np.abs(dup-orig)/(dup+orig)*100; sub = sub.assign(HARD=hard)

#             fig, ax = plt.subplots(1,2,figsize=(14,5))
#             # HARD %
#             sort = sub.sort_values("HARD"); sort["pct"]=sort["HARD"].rank(pct=True)*100
#             lim = hard_pct.get(dtyp,.3)*100
#             ax[0].plot(sort["pct"], sort["HARD"],lw=.6,color="#2F4F4F")
#             ax[0].scatter(sort["pct"], sort["HARD"],s=8,color="black")
#             ax[0].axhline(lim,ls="--",color="#86A586")
#             ax[0].set_ylim(0,100); ax[0].set_title(f"HARD – {el}")
#             ax[0].set_xlabel("Percentile"); ax[0].set_ylabel("HARD (%)")
#             total=len(sub); fails=(sub["HARD"]>lim).sum(); fr=fails/total*100
#             ax[0].text(.02,.98,f"Pairs {total}\nFails {fails}\nRate {fr:.1f}%",
#                        transform=ax[0].transAxes,va="top",
#                        bbox=dict(facecolor="white",alpha=.7))

#             # Scatter
#             ax[1].scatter(orig, dup, s=10, color="black")
#             if use_log: ax[1].set_xscale("log"); ax[1].set_yscale("log")
#             if custom_lim:
#                 ax[1].set_xlim(x_min, x_max); ax[1].set_ylim(y_min, y_max)
#             else:
#                 hi = max(orig.max(), dup.max())*1.05; ax[1].set_xlim(0,hi); ax[1].set_ylim(0,hi)
#             pct = hard_pct.get(dtyp,.3)
#             ax[1].plot(ax[1].get_xlim(), ax[1].get_ylim(), lw=1, color="#C1E27C")
#             ax[1].plot([0,ax[1].get_xlim()[1]], [0, ax[1].get_ylim()[1]*(1+pct)], ls="--",color="#86A586")
#             ax[1].plot([0,ax[1].get_xlim()[1]], [0, ax[1].get_ylim()[1]*(1-pct)], ls="--",color="#86A586")
#             ax[1].set_title(f"{lab} • {dtyp}"); ax[1].set_xlabel("Original"); ax[1].set_ylabel("Duplicate")

#             st.pyplot(fig)
#             ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(
#                 tmp_fig(fig), Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))

#             summary_rows.append({"Lab":lab,"DupType":dtyp,"Element":el,"Pairs":total,
#                                  "Fails":fails,"FailRate%":round(fr,2),
#                                  "Corr":round(orig.corr(dup),3),
#                                  "MeanΔ":round((dup-orig).mean(),2)})

# # ── 10. Summary + downloads ─────────────────────────────────────────
# summary = pd.DataFrame(summary_rows)
# st.subheader("📊 Summary"); st.dataframe(summary)
# xls = io.BytesIO(); summary.to_excel(xls,index=False)
# ppt_bytes = io.BytesIO(); ppt.save(ppt_bytes)
# st.download_button("Download summary.xlsx", xls.getvalue(), file_name="summary_duplicates.xlsx")
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="duplicate_plots.pptx")

###### Version 2 ##############

# import os, io, tempfile, re
# import numpy as np
# import pandas as pd
# import matplotlib.pyplot as plt
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches

# st.set_page_config(page_title="Duplicate QAQC Charts", layout="wide")
# st.title("🪓 Duplicate QAQC Charts")

# # ──────────── helpers ────────────
# _det_pat  = re.compile(r"^[<>]=?|^[≤≥]")
# _unit_pat = re.compile(r"\s*\(.*\)\s*$")
# def clean_numeric(s):
#     if pd.api.types.is_numeric_dtype(s):
#         return pd.to_numeric(s, errors="coerce")
#     s = (
#         s.astype(str)
#         .str.strip()
#         .str.replace(_unit_pat, "", regex=True)
#         .str.replace(_det_pat,  "", regex=True)
#         .str.replace(",", "", regex=False)
#         .replace({"": np.nan, "nan": np.nan, "None": np.nan, "NULL": np.nan})
#     )
#     return pd.to_numeric(s, errors="coerce")

# def tmp_fig(fig):
#     f = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(f.name, dpi=300, bbox_inches="tight")
#     plt.close(fig); return f.name

# # ──────────── upload ────────────
# csv = st.sidebar.file_uploader("Upload duplicate CSV", ["csv"])
# if not csv: st.stop()
# df = pd.read_csv(csv)

# # ───── mandatory column mapping (stacked format) ─────
# with st.sidebar.expander("🗺️  Column mapping", True):
#     lab_col = st.selectbox("Lab column", df.columns)
#     dup_col = st.selectbox("Duplicate‑type column", df.columns)

#     # *** mandatory element column (no 'None') ***
#     element_col = st.selectbox("Element column", df.columns)

#     orig_val_col = st.selectbox("Original value column", df.columns)
#     dup_val_col  = st.selectbox("Duplicate value column",  df.columns)

# # ───── filters ─────
# with st.sidebar.expander("🔎 Core filters"):
#     labs_sel = st.multiselect("Labs", df[lab_col].dropna().unique(), default=list(df[lab_col].dropna().unique()))
#     dup_sel  = st.multiselect("Duplicate types", df[dup_col].dropna().unique(), default=list(df[dup_col].dropna().unique()))

# cat_filters = []
# with st.sidebar.expander("🎛️ Extra categorical filters"):
#     for i in range(1,4):
#         col = st.selectbox(f"Filter {i}", ["None"]+list(df.columns), key=f"cat_{i}")
#         if col != "None":
#             cat_filters.append((col, st.multiselect("Values", df[col].dropna().unique(), key=f"vals_{i}")))

# with st.sidebar.expander("🗓️ Date filter"):
#     date_col = st.selectbox("Date column", ["None"]+list(df.columns), index=0)
#     date_rng = None
#     if date_col != "None":
#         df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
#         dmin, dmax = map(lambda d: d.date() if pd.notna(d) else None, (df[date_col].min(), df[date_col].max()))
#         date_rng = st.date_input("Range", (dmin, dmax))

# # ───── style ─────
# with st.sidebar.expander("📐 Plot options & style"):
#     use_log    = st.checkbox("Log‑scale axes (scatter)")
#     custom_lim = st.checkbox("Custom X/Y limits")
#     x_min=y_min=x_max=y_max=None
#     if custom_lim:
#         x_min=st.number_input("X min", value=0.0); x_max=st.number_input("X max", value=0.0)
#         y_min=st.number_input("Y min", value=0.0); y_max=st.number_input("Y max", value=0.0)

#     st.markdown("**Global style**")
#     line_width = st.number_input("Line width (pt)", min_value=0.1, value=0.8, step=0.1, format="%.2f")
#     point_size = st.number_input("Point size",      min_value=1,   value=8,   step=1)
#     main_line_col  = st.color_picker("Curve / identity line colour", "#2F4F4F")
#     point_col      = st.color_picker("Point colour", "#000000")
#     limit_line_col = st.color_picker("Limit / threshold line colour", "#86A586")

# # ───── element list ─────
# unique_elems = sorted(df[element_col].dropna().unique())
# sel_elems    = st.sidebar.multiselect("Elements", unique_elems, default=unique_elems)
# if not st.sidebar.button("🚀 Run") or not sel_elems:
#     st.stop()

# # ───── constants ─────
# ppt, summary = Presentation(), []
# hard_pct = {"PD":0.1, "CD":0.2, "FD":0.3}

# # ───── main loop ─────
# for el in sel_elems:
#     data = df[df[element_col]==el]
#     data = data[data[lab_col].isin(labs_sel) & data[dup_col].isin(dup_sel)]
#     for c,v in cat_filters: data = data[data[c].isin(v)] if v else data
#     if date_col!="None" and date_rng:
#         start,end = map(pd.to_datetime, date_rng); end += pd.Timedelta(days=1)
#         data = data[(data[date_col]>=start)&(data[date_col]<end)]
#     if data.empty:
#         st.warning(f"No data for {el}")
#         continue

#     for lab in data[lab_col].unique():
#         for dtyp in data[dup_col].unique():
#             sub = data[(data[lab_col]==lab)&(data[dup_col]==dtyp)]
#             orig = clean_numeric(sub[orig_val_col]); dup = clean_numeric(sub[dup_val_col])
#             m = orig.notna() & dup.notna(); sub,orig,dup = sub[m].copy(), orig[m], dup[m]
#             if orig.empty: continue

#             denom = dup+orig
#             sub["HARD"] = np.where(denom>0, np.abs(dup-orig)/denom*100, np.nan)
#             sub.dropna(subset=["HARD"], inplace=True)
#             if sub.empty: continue
#             orig,dup = orig.loc[sub.index], dup.loc[sub.index]

#             fig, ax = plt.subplots(1,2,figsize=(14,5))

#             # HARD plot ----------------------------------------------------------------
#             sort=sub.sort_values("HARD"); sort["pct"]=sort["HARD"].rank(pct=True)*100
#             lim=hard_pct.get(dtyp,0.3)*100
#             ax[0].plot(sort["pct"], sort["HARD"], lw=line_width, color=main_line_col)
#             ax[0].scatter(sort["pct"], sort["HARD"], s=point_size, color=point_col)
#             ax[0].axhline(lim, ls="--", lw=line_width, color=limit_line_col)
#             ax[0].set_ylim(0,100); ax[0].set_xlabel("Percentile"); ax[0].set_ylabel("HARD (%)")
#             ax[0].set_title(f"HARD index for {el} - {dtyp} - {lab}")
#             total=len(sub); fails=(sub["HARD"]>lim).sum(); fr=fails/total*100
#             ax[0].text(.02,.98,f"Pairs: {total}\nFails: {fails}\nRate: {fr:.1f}%",
#                        transform=ax[0].transAxes, va="top", bbox=dict(fc="white",alpha=.7))

#             # Scatter plot -------------------------------------------------------------
#             ax[1].scatter(orig, dup, s=point_size, color=point_col, label="Pairs")
#             if use_log: ax[1].set_xscale("log"); ax[1].set_yscale("log")
#             if custom_lim: ax[1].set_xlim(x_min,x_max); ax[1].set_ylim(y_min,y_max)
#             else:
#                 hi=max(orig.max(),dup.max())*1.05; ax[1].set_xlim(0,hi); ax[1].set_ylim(0,hi)
#             pct = hard_pct.get(dtyp,0.3)
#             x_max_lim=ax[1].get_xlim()[1]; y_max_lim=ax[1].get_ylim()[1]
#             ax[1].plot([0,x_max_lim],[0,y_max_lim],lw=line_width,color=main_line_col,label="y=x")
#             ax[1].plot([0,x_max_lim],[0,y_max_lim*(1+pct)],ls="--",lw=line_width,color=limit_line_col,
#                        label=f"+{pct*100:.0f}%")
#             ax[1].plot([0,x_max_lim],[0,y_max_lim*(1-pct)],ls="--",lw=line_width,color=limit_line_col,
#                        label=f"-{pct*100:.0f}%")
#             ax[1].set_title(f"Scatter Plot for {el} - {dtyp} - {lab}")
#             ax[1].set_xlabel("Original"); ax[1].set_ylabel("Duplicate")
#             ax[1].legend(loc="upper left")

#             st.pyplot(fig)
#             ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(tmp_fig(fig), Inches(1), Inches(1),
#                                                                            width=Inches(8), height=Inches(4.5))

#             summary.append({"Lab":lab,"DupType":dtyp,"Element":el,
#                             "Pairs":total,"Fails":fails,"FailRate%":round(fr,2),
#                             "Corr":round(orig.corr(dup),3) if len(orig)>1 else None,
#                             "MeanΔ":round((dup-orig).mean(),2)})

# # ───── summary & downloads ─────
# sum_df = pd.DataFrame(summary)
# st.subheader("📊 Summary"); st.dataframe(sum_df)

# xls, ppt_bytes = io.BytesIO(), io.BytesIO()
# sum_df.to_excel(xls, index=False); ppt.save(ppt_bytes)

# st.download_button("Download summary.xlsx", xls.getvalue(), file_name="summary_duplicates.xlsx")
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="duplicate_plots.pptx")

##### Version 3 ########

# import os, io, json, tempfile, re
# import numpy as np
# import pandas as pd
# import matplotlib.pyplot as plt
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches

# st.set_page_config(page_title="Duplicate QAQC Charts", layout="wide")
# st.title("🪓 Duplicate QAQC Charts")

# # ───────────────────────────── helpers ──────────────────────────────
# _det_pat  = re.compile(r"^[<>]=?|^[≤≥]")
# _unit_pat = re.compile(r"\s*\(.*\)\s*$")
# def clean_numeric(s):
#     if pd.api.types.is_numeric_dtype(s):
#         return pd.to_numeric(s, errors="coerce")
#     s = (
#         s.astype(str)
#         .str.strip()
#         .str.replace(_unit_pat, "", regex=True)
#         .str.replace(_det_pat,  "", regex=True)
#         .str.replace(",", "", regex=False)
#         .replace({"": np.nan, "nan": np.nan, "None": np.nan, "NULL": np.nan})
#     )
#     return pd.to_numeric(s, errors="coerce")

# def tmp_fig(fig):
#     f = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(f.name, dpi=300, bbox_inches="tight")
#     plt.close(fig)
#     return f.name

# def idx_of(val, options):
#     """Return index of *val* in *options*, else 0."""
#     try:
#         return options.index(val)
#     except ValueError:
#         return 0

# def safe_defaults(defaults, options):
#     """Filter a default list so only existing options remain."""
#     return [d for d in (defaults or []) if d in options]

# # ─────────────────────── 1) optional config load ────────────────────
# if "loaded_cfg" not in st.session_state:
#     st.session_state.loaded_cfg = False

# cfg_file = st.sidebar.file_uploader("🔄 Load saved config (JSON)", type=["json"])
# if cfg_file and not st.session_state.loaded_cfg:
#     st.session_state.update(json.load(cfg_file))
#     st.session_state.loaded_cfg = True
#     st.success("Config loaded!")

# # ───────────────────────── 2) load CSV ──────────────────────────────
# csv_bytes = st.sidebar.file_uploader("📂 Upload CSV", type=["csv"], key="csv_bytes")

# df = None
# if csv_bytes is not None:
#     df = pd.read_csv(csv_bytes)
#     st.session_state["csv_path"] = ""        # clear stored path if user uploads
# elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
#     df = pd.read_csv(st.session_state["csv_path"])

# if df is None:
#     st.info("Upload a CSV (or load a config with a valid path).")
#     st.stop()

# # ─────────────────────── 3) column mapping ──────────────────────────
# with st.sidebar.expander("🗺️ Column mapping", True):
#     lab_col      = st.selectbox("Lab column", df.columns,
#                                 index=idx_of(st.session_state.get("lab_col", df.columns[0]), list(df.columns)),
#                                 key="lab_col")
#     dup_col      = st.selectbox("Duplicate‑type column", df.columns,
#                                 index=idx_of(st.session_state.get("dup_col", df.columns[0]), list(df.columns)),
#                                 key="dup_col")
#     element_col  = st.selectbox("Element column (required)", df.columns,
#                                 index=idx_of(st.session_state.get("element_col", df.columns[0]), list(df.columns)),
#                                 key="element_col")
#     orig_val_col = st.selectbox("Original value column", df.columns,
#                                 index=idx_of(st.session_state.get("orig_val_col", df.columns[0]), list(df.columns)),
#                                 key="orig_val_col")
#     dup_val_col  = st.selectbox("Duplicate value column", df.columns,
#                                 index=idx_of(st.session_state.get("dup_val_col", df.columns[0]), list(df.columns)),
#                                 key="dup_val_col")

# # ───────────────────────── 4) filters ───────────────────────────────
# with st.sidebar.expander("🔎 Core filters"):
#     labs_opts = list(df[lab_col].dropna().unique())
#     dup_opts  = list(df[dup_col].dropna().unique())
#     labs_sel = st.multiselect("Labs", labs_opts,
#                               default=safe_defaults(st.session_state.get("labs_sel"), labs_opts),
#                               key="labs_sel")
#     dup_sel  = st.multiselect("Duplicate types", dup_opts,
#                               default=safe_defaults(st.session_state.get("dup_sel"), dup_opts),
#                               key="dup_sel")

# cat_filters = []
# with st.sidebar.expander("🎛️ Extra categorical filters"):
#     for i in range(1, 3 + 1):
#         col_key, vals_key = f"cat_col_{i}", f"cat_vals_{i}"
#         col_opts = ["None"] + list(df.columns)
#         col = st.selectbox(f"Filter {i} column", col_opts,
#                            index=idx_of(st.session_state.get(col_key, "None"), col_opts),
#                            key=col_key)
#         if col != "None":
#             vals_opts = list(df[col].dropna().unique())
#             vals_sel  = st.multiselect("Values", vals_opts,
#                                        default=safe_defaults(st.session_state.get(vals_key), vals_opts),
#                                        key=vals_key)
#             cat_filters.append((col, vals_sel))

# with st.sidebar.expander("🗓️ Date filter"):
#     date_opts = ["None"] + list(df.columns)
#     date_col  = st.selectbox("Date column", date_opts,
#                              index=idx_of(st.session_state.get("date_col", "None"), date_opts),
#                              key="date_col")
#     date_rng  = None
#     if date_col != "None":
#         df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
#         dmin, dmax = map(lambda d: d.date() if pd.notna(d) else None,
#                          (df[date_col].min(), df[date_col].max()))
#         date_rng = st.date_input("Range",
#                                  st.session_state.get("date_range", (dmin, dmax)),
#                                  key="date_range")

# # ───────────────────────── 5) style ─────────────────────────────────
# with st.sidebar.expander("📐 Plot options & style"):
#     use_log    = st.checkbox("Log‑scale axes (scatter)",
#                              value=st.session_state.get("use_log", False),
#                              key="use_log")
#     custom_lim = st.checkbox("Custom X/Y limits",
#                              value=st.session_state.get("custom_lim", False),
#                              key="custom_lim")
#     if custom_lim:
#         x_min = st.number_input("X min", value=st.session_state.get("x_min", 0.0), key="x_min")
#         x_max = st.number_input("X max", value=st.session_state.get("x_max", 0.0), key="x_max")
#         y_min = st.number_input("Y min", value=st.session_state.get("y_min", 0.0), key="y_min")
#         y_max = st.number_input("Y max", value=st.session_state.get("y_max", 0.0), key="y_max")
#     else:
#         x_min = y_min = x_max = y_max = None

#     st.markdown("**Global style**")
#     line_width = st.number_input("Line width (pt)", min_value=0.1,
#                                  value=st.session_state.get("line_width", 0.8),
#                                  step=0.1, format="%.2f", key="line_width")
#     point_size = st.number_input("Point size", min_value=1,
#                                  value=st.session_state.get("point_size", 8),
#                                  step=1, key="point_size")
#     main_line_col  = st.color_picker("Curve / identity line colour",
#                                      st.session_state.get("main_line_col", "#2F4F4F"),
#                                      key="main_line_col")
#     point_col      = st.color_picker("Point colour",
#                                      st.session_state.get("point_col", "#000000"),
#                                      key="point_col")
#     limit_line_col = st.color_picker("Limit / threshold line colour",
#                                      st.session_state.get("limit_line_col", "#86A586"),
#                                      key="limit_line_col")

# # ───────────────────────── 6) element selection ─────────────────────
# unique_elems = sorted(df[element_col].dropna().unique())
# sel_elems = st.sidebar.multiselect("Elements", unique_elems,
#                                    default=safe_defaults(st.session_state.get("sel_elems"), unique_elems),
#                                    key="sel_elems")

# # ───────────────────────── 7) run flag ──────────────────────────────
# run = st.session_state.loaded_cfg or st.sidebar.button("🚀 Run")
# if not run:
#     st.stop()

# # ───────────────────────── 8) analysis ──────────────────────────────
# hard_pct = {"PD": 0.1, "CD": 0.2, "FD": 0.3}
# ppt, summary = Presentation(), []

# for el in sel_elems:
#     data = df[df[element_col] == el]
#     data = data[data[lab_col].isin(labs_sel) & data[dup_col].isin(dup_sel)]
#     for c, v in cat_filters:
#         data = data[data[c].isin(v)] if v else data
#     if date_col != "None" and date_rng:
#         start, end = map(pd.to_datetime, date_rng)
#         end += pd.Timedelta(days=1)
#         data = data[(data[date_col] >= start) & (data[date_col] < end)]
#     if data.empty:
#         st.warning(f"No data for {el}")
#         continue

#     for lab in data[lab_col].unique():
#         for dtyp in data[dup_col].unique():
#             sub = data[(data[lab_col] == lab) & (data[dup_col] == dtyp)]
#             orig = clean_numeric(sub[orig_val_col])
#             dup  = clean_numeric(sub[dup_val_col])
#             m = orig.notna() & dup.notna()
#             sub, orig, dup = sub[m].copy(), orig[m], dup[m]
#             if orig.empty:
#                 continue

#             denom = dup + orig
#             sub["HARD"] = np.where(denom > 0, np.abs(dup - orig) / denom * 100, np.nan)
#             sub.dropna(subset=["HARD"], inplace=True)
#             if sub.empty:
#                 continue
#             orig, dup = orig.loc[sub.index], dup.loc[sub.index]

#             fig, ax = plt.subplots(1, 2, figsize=(14, 5))

#             # HARD plot --------------------------------------------------
#             sort = sub.sort_values("HARD")
#             sort["pct"] = sort["HARD"].rank(pct=True) * 100
#             lim = hard_pct.get(dtyp, 0.3) * 100
#             ax[0].plot(sort["pct"], sort["HARD"], lw=line_width, color=main_line_col)
#             ax[0].scatter(sort["pct"], sort["HARD"], s=point_size, color=point_col)
#             ax[0].axhline(lim, ls="--", lw=line_width, color=limit_line_col)
#             ax[0].set_ylim(0, 100)
#             ax[0].set_xlabel("Percentile")
#             ax[0].set_ylabel("HARD (%)")
#             ax[0].set_title(f"HARD index for {el} - {dtyp} - {lab}")
#             total = len(sub)
#             fails = (sub["HARD"] > lim).sum()
#             fr = fails / total * 100
#             ax[0].text(
#                 0.02,
#                 0.98,
#                 f"Pairs: {total}\nFails: {fails}\nRate: {fr:.1f}%",
#                 transform=ax[0].transAxes,
#                 va="top",
#                 bbox=dict(fc="white", alpha=0.7),
#             )

#             # Scatter plot -----------------------------------------------
#             ax[1].scatter(orig, dup, s=point_size, color=point_col, label="Pairs")
#             if use_log:
#                 ax[1].set_xscale("log")
#                 ax[1].set_yscale("log")
#             if custom_lim:
#                 ax[1].set_xlim(x_min, x_max)
#                 ax[1].set_ylim(y_min, y_max)
#             else:
#                 hi = max(orig.max(), dup.max()) * 1.05
#                 ax[1].set_xlim(0, hi)
#                 ax[1].set_ylim(0, hi)
#             pct = hard_pct.get(dtyp, 0.3)
#             x_max_lim = ax[1].get_xlim()[1]
#             y_max_lim = ax[1].get_ylim()[1]
#             ax[1].plot([0, x_max_lim], [0, y_max_lim], lw=line_width, color=main_line_col, label="y=x")
#             ax[1].plot([0, x_max_lim], [0, y_max_lim * (1 + pct)], ls="--", lw=line_width,
#                        color=limit_line_col, label=f"+{pct*100:.0f}%")
#             ax[1].plot([0, x_max_lim], [0, y_max_lim * (1 - pct)], ls="--", lw=line_width,
#                        color=limit_line_col, label=f"-{pct*100:.0f}%")
#             ax[1].set_title(f"Scatter Plot for {el} - {dtyp} - {lab}")
#             ax[1].set_xlabel("Original")
#             ax[1].set_ylabel("Duplicate")
#             ax[1].legend(loc="upper left")

#             st.pyplot(fig)
#             ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(
#                 tmp_fig(fig), Inches(1), Inches(1), width=Inches(8), height=Inches(4.5)
#             )

#             summary.append(
#                 {
#                     "Lab": lab,
#                     "DupType": dtyp,
#                     "Element": el,
#                     "Pairs": total,
#                     "Fails": fails,
#                     "FailRate%": round(fr, 2),
#                     "Corr": round(orig.corr(dup), 3) if len(orig) > 1 else None,
#                     "MeanΔ": round((dup - orig).mean(), 2),
#                 }
#             )

# # ─────────────────────── summary & downloads ────────────────────────
# sum_df = pd.DataFrame(summary)
# st.subheader("📊 Summary")
# st.dataframe(sum_df)

# xls = io.BytesIO(); sum_df.to_excel(xls, index=False)
# ppt_bytes = io.BytesIO(); ppt.save(ppt_bytes)

# st.download_button("Download summary.xlsx", xls.getvalue(), file_name="summary_duplicates.xlsx")
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="duplicate_plots.pptx")

# # ─────────────────────── save config (skip blobs) ────────────────────
# serialisable = {}
# for k, v in st.session_state.items():
#     if k.startswith("_"):
#         continue
#     try:
#         json.dumps(v)
#         serialisable[k] = v
#     except TypeError:
#         # skip non‑serialisable objects like UploadedFile
#         pass

# cfg_bytes = io.BytesIO(json.dumps(serialisable, indent=2).encode("utf-8"))

# st.sidebar.download_button("💾 Download config JSON",
#                            cfg_bytes.getvalue(),
#                            file_name="qaqc_config.json",
#                            mime="application/json")


##### Version 4 ########
import os, io, json, tempfile, re
import numpy as np, pandas as pd, matplotlib.pyplot as plt, streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pyrpa.UI import common

st.set_page_config(page_title="Duplicate QAQC Charts", layout="wide")
st.title("🪓 Duplicate QAQC Charts")

# ───────────────────────── helpers ─────────────────────────
_det_pat = re.compile(r"^[<>]=?|^[≤≥]")
def clean_numeric(s):
    if pd.api.types.is_numeric_dtype(s):
        return pd.to_numeric(s, errors="coerce")
    s = (
        s.astype(str)
        .str.strip()
        .str.replace(_det_pat, "", regex=True)
        .str.replace(",", "", regex=False)
        .str.replace(r"[^\d\.\-eE+]", "", regex=True)
    )
    return pd.to_numeric(s, errors="coerce")

def tmp_fig(fig):
    t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(t.name, dpi=300, bbox_inches="tight"); plt.close(fig)
    return t.name

def idx_of(val, opts): return opts.index(val) if val in opts else 0
def safe_defaults(defs, opts): return [d for d in (defs or []) if d in opts]

# ───────── 1) optional config load ─────────
if "cfg_loaded" not in st.session_state:
    st.session_state.cfg_loaded = False

cfg_file = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
if cfg_file and not st.session_state.cfg_loaded:
    st.session_state.update(json.load(cfg_file))
    st.session_state.cfg_loaded = True
    st.success("Config loaded!")

# ───────── 2) data load ─────────
csv_bytes = st.sidebar.file_uploader("📂 Upload data (CSV or Excel)", ["csv","xlsx","xls"], key="csv_bytes")
df = None
if csv_bytes is not None:
    df = common.read_data_file(csv_bytes)
    st.session_state["csv_path"] = ""
elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
    df = common.read_data_file(st.session_state["csv_path"])

if df is None:
    st.info("Upload CSV or load config."); st.stop()

# ───────── 3) column mapping ─────────
with st.sidebar.expander("🗺️ Column mapping", True):
    lab_col      = st.selectbox("Lab column", df.columns,
                                index=idx_of(st.session_state.get("lab_col", df.columns[0]), list(df.columns)),
                                key="lab_col")
    dup_col      = st.selectbox("Duplicate‑type column", df.columns,
                                index=idx_of(st.session_state.get("dup_col", df.columns[0]), list(df.columns)),
                                key="dup_col")
    element_col  = st.selectbox("Element column", df.columns,
                                index=idx_of(st.session_state.get("element_col", df.columns[0]), list(df.columns)),
                                key="element_col")
    unit_col     = st.selectbox("Unit column", df.columns,
                                index=idx_of(st.session_state.get("unit_col", df.columns[0]), list(df.columns)),
                                key="unit_col")
    orig_val_col = st.selectbox("Original value column", df.columns,
                                index=idx_of(st.session_state.get("orig_val_col", df.columns[0]), list(df.columns)),
                                key="orig_val_col")
    dup_val_col  = st.selectbox("Duplicate value column", df.columns,
                                index=idx_of(st.session_state.get("dup_val_col", df.columns[0]), list(df.columns)),
                                key="dup_val_col")

# ───────── 4) filters ─────────
with st.sidebar.expander("🔎 Core filters"):
    labs_sel = st.multiselect(
        "Labs",
        df[lab_col].dropna().unique(),
        default=safe_defaults(st.session_state.get("labs_sel"), df[lab_col].unique()),
        key="labs_sel",
    )
    dup_sel = st.multiselect(
        "Duplicate types",
        df[dup_col].dropna().unique(),
        default=safe_defaults(st.session_state.get("dup_sel"), df[dup_col].unique()),
        key="dup_sel",
    )

cat_filters, num_filters = [], []

with st.sidebar.expander("🎛️ Categorical filters"):
    for i in range(1, 4):
        col_key, vals_key = f"cat_col_{i}", f"cat_vals_{i}"
        col = st.selectbox(
            f"Cat {i} column",
            ["None"] + list(df.columns),
            index=idx_of(st.session_state.get(col_key, "None"), ["None"] + list(df.columns)),
            key=col_key,
        )
        if col != "None":
            vals_sel = st.multiselect(
                "Values",
                df[col].dropna().unique(),
                default=safe_defaults(st.session_state.get(vals_key), df[col].unique()),
                key=vals_key,
            )
            cat_filters.append((col, vals_sel))

with st.sidebar.expander("🔢 Numerical filters"):
    num_cols = [c for c in df.columns if clean_numeric(df[c]).notna().any()]
    for i in range(1, 3):
        col_key, rng_key = f"num_col_{i}", f"num_rng_{i}"
        col = st.selectbox(
            f"Num {i} column",
            ["None"] + num_cols,
            index=idx_of(st.session_state.get(col_key, "None"), ["None"] + num_cols),
            key=col_key,
        )
        if col != "None":
            s = clean_numeric(df[col]); lo, hi = float(s.min()), float(s.max())
            rng = st.slider(
                "Range", lo, hi,
                value=st.session_state.get(rng_key, (lo, hi)),
                key=rng_key,
            )
            num_filters.append((col, rng))

with st.sidebar.expander("🗓️ Date filter"):
    date_col = st.selectbox(
        "Date column",
        ["None"] + list(df.columns),
        index=idx_of(st.session_state.get("date_col", "None"), ["None"] + list(df.columns)),
        key="date_col",
    )
    date_rng = None
    if date_col != "None":
        df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
        dmin, dmax = df[date_col].min(), df[date_col].max()
        date_rng = st.date_input(
            "Range",
            st.session_state.get(
                "date_range",
                (
                    None if pd.isna(dmin) else dmin.date(),
                    None if pd.isna(dmax) else dmax.date(),
                ),
            ),
            key="date_range",
        )

# ───────── 5) style & Nugget effect ─────────
with st.sidebar.expander("📐 Style & Nugget"):
    line_width = st.number_input(
        "Line width", 0.1, value=float(st.session_state.get("line_width", 0.8)),
        step=0.1, format="%.2f", key="line_width"
    )
    point_size = st.number_input(
        "Point size", 1, value=int(st.session_state.get("point_size", 8)),
        step=1, key="point_size"
    )
    main_line_col  = st.color_picker("Curve / identity", st.session_state.get("main_line_col", "#2F4F4F"), key="main_line_col")
    point_col      = st.color_picker("Point colour",    st.session_state.get("point_col", "#000000"), key="point_col")
    limit_line_col = st.color_picker("Limit colour",    st.session_state.get("limit_line_col", "#86A586"), key="limit_line_col")
    nugget_on  = st.checkbox("Nugget effect (+10 pp)", value=st.session_state.get("nugget_on", False), key="nugget_on")
    nugget_col = st.color_picker("Nugget line colour",  st.session_state.get("nugget_col", "#FF0000"), key="nugget_col")
    use_log   = st.checkbox("Log axes (scatter)", value=st.session_state.get("use_log", False), key="use_log")
    custom_lim = st.checkbox("Custom axis limits", value=st.session_state.get("custom_lim", False), key="custom_lim")
    x_min=y_min=x_max=y_max=None
    if custom_lim:
        x_min = st.number_input("X min", value=st.session_state.get("x_min", 0.0), key="x_min")
        x_max = st.number_input("X max", value=st.session_state.get("x_max", 0.0), key="x_max")
        y_min = st.number_input("Y min", value=st.session_state.get("y_min", 0.0), key="y_min")
        y_max = st.number_input("Y max", value=st.session_state.get("y_max", 0.0), key="y_max")

# ───────── 6) element selection ─────────
unique_elems = sorted(df[element_col].dropna().unique())
sel_elems = st.sidebar.multiselect(
    "Elements", unique_elems,
    default=safe_defaults(st.session_state.get("sel_elems"), unique_elems),
    key="sel_elems"
)

# ───────── 7) run flag ─────────
if not (st.session_state.cfg_loaded or st.sidebar.button("🚀 Run")):
    st.stop()

# ───────── 8) analysis ─────────
hard_base = {"PD": 0.1, "CD": 0.2, "FD": 0.3}
ppt, summary_rows = Presentation(), []

for el in sel_elems:
    data = df[df[element_col] == el]
    data = data[data[lab_col].isin(labs_sel) & data[dup_col].isin(dup_sel)]

    for c, v in cat_filters: data = data[data[c].isin(v)] if v else data
    for c, (lo, hi) in num_filters:
        num = clean_numeric(data[c]); data = data[(num >= lo) & (num <= hi)]

    if date_col != "None" and date_rng:
        start, end = map(pd.to_datetime, date_rng); end += pd.Timedelta(days=1)
        data = data[(data[date_col] >= start) & (data[date_col] < end)]

    if data.empty: continue

    for lab in data[lab_col].unique():
        for dtyp in data[dup_col].unique():
            sub = data[(data[lab_col] == lab) & (data[dup_col] == dtyp)]
            orig, dup = clean_numeric(sub[orig_val_col]), clean_numeric(sub[dup_val_col])
            mask = orig.notna() & dup.notna()
            sub, orig, dup = sub[mask].copy(), orig[mask], dup[mask]
            if orig.empty: continue

            sub["HARD"] = np.abs(dup - orig) / (dup + orig) * 100
            if sub.empty: continue

            # -------- titles: get unit --------------
            unit_val = ""
            if unit_col in sub.columns:
                uvals = sub[unit_col].dropna().unique()
                if len(uvals): unit_val = f"{uvals[0]}"

            pct_base = hard_base.get(dtyp, 0.3)
            pct_nug  = pct_base + 0.1 if nugget_on else None
            lim_base = pct_base * 100
            lim_nug  = pct_nug  * 100 if nugget_on else None

            # -------- figure ----------
            fig, ax = plt.subplots(1, 2, figsize=(14, 5))

            # HARD plot (no legend)
            srt = sub.sort_values("HARD")
            srt["pct"] = srt["HARD"].rank(pct=True) * 100
            ax[0].plot(srt["pct"], srt["HARD"], lw=line_width, color=main_line_col)
            ax[0].scatter(srt["pct"], srt["HARD"], s=point_size, color=point_col)
            ax[0].axhline(lim_base, ls="--", lw=line_width, color=limit_line_col)
            if nugget_on:
                ax[0].axhline(lim_nug, ls="--", lw=line_width, color=nugget_col)
            ax[0].set_ylim(0, 100)
            ax[0].set_xlabel("Percentile"); ax[0].set_ylabel("HARD (%)")
            ax[0].set_title(f"HARD index for {el} {unit_val} - {dtyp} - {lab}".strip())
            total = len(sub)
            fails_base = (sub["HARD"] > lim_base).sum()
            rate_base = fails_base / total * 100

            if nugget_on:
                fails_nug = (sub["HARD"] > lim_nug).sum()
                rate_nug = fails_nug / total * 100
                stats_text = (
                    f"Pairs: {total}\n"
                    f"Fails (nugget): {fails_nug} (Rate: {rate_nug:.1f}%)\n"
                    f"Fails (base): {fails_base} (Rate: {rate_base:.1f}%)"
                )

                summary_fails = fails_nug
                summary_rate = rate_nug
            else:
                stats_text = f"Pairs: {total}\nFails: {fails_base}\nRate: {rate_base:.1f}%"

                summary_fails = fails_base
                summary_rate = rate_base

            ax[0].text(
                0.02, 0.98,
                stats_text,
                transform=ax[0].transAxes,
                va="top",
                bbox=dict(fc="white", alpha=0.7),
            )

            # Scatter plot (legend kept)
            ax[1].scatter(orig, dup, s=point_size, color=point_col, label="Pairs")
            if use_log:
                ax[1].set_xscale("log"); ax[1].set_yscale("log")
            if custom_lim:
                ax[1].set_xlim(x_min, x_max); ax[1].set_ylim(y_min, y_max)
            else:
                hi = max(orig.max(), dup.max()) * 1.05; ax[1].set_xlim(0, hi); ax[1].set_ylim(0, hi)
            xlim, ylim = ax[1].get_xlim()[1], ax[1].get_ylim()[1]
            ax[1].plot([0, xlim], [0, ylim], lw=line_width, color=main_line_col, label="y=x")
            ax[1].plot([0, xlim], [0, ylim * (1 + pct_base)], ls="--", lw=line_width, color=limit_line_col,
                       label=f"+{pct_base * 100:.0f}%")
            ax[1].plot([0, xlim], [0, ylim * (1 - pct_base)], ls="--", lw=line_width, color=limit_line_col,
                       label=f"-{pct_base * 100:.0f}%")
            if nugget_on:
                ax[1].plot([0, xlim], [0, ylim * (1 + pct_nug)], ls="--", lw=line_width, color=nugget_col,
                           label=f"+{pct_nug * 100:.0f}%")
                ax[1].plot([0, xlim], [0, ylim * (1 - pct_nug)], ls="--", lw=line_width, color=nugget_col,
                           label=f"-{pct_nug * 100:.0f}%")
            ax[1].set_title(f"Scatter Plot for {el} {unit_val} - {dtyp} - {lab}".strip())
            ax[1].set_xlabel("Original"); ax[1].set_ylabel("Duplicate"); ax[1].legend(loc="upper left")

            st.pyplot(fig)
            ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(
                tmp_fig(fig), Inches(1), Inches(1), width=Inches(8), height=Inches(4.5)
            )

            summary_rows.append({
                "Lab": lab, "DupType": dtyp, "Element": el,
                "Pairs": total, "Fails": summary_fails, "FailRate%": round(summary_rate, 2),
                "Corr": round(orig.corr(dup), 3) if len(orig) > 1 else None,
                "MeanΔ": round((dup - orig).mean(), 2),
                "Min": round(orig.min(), 2), "Max": round(orig.max(), 2), "Mean": round(orig.mean(), 2)
            })

# ───────── summary / export ─────────
summary_df = pd.DataFrame(summary_rows)
st.subheader("📊 Summary"); st.dataframe(summary_df)

xls        = io.BytesIO(); summary_df.to_excel(xls, index=False)
ppt_bytes  = io.BytesIO(); ppt.save(ppt_bytes)

st.download_button("Download summary.xlsx", xls.getvalue(), file_name="summary_duplicates.xlsx")
st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="duplicate_plots.pptx")

# ───────── save config (skip blobs) ─────────
serialisable={}
for k,v in st.session_state.items():
    if k.startswith("_"): continue
    try: json.dumps(v); serialisable[k]=v
    except TypeError: pass

cfg_bytes=io.BytesIO(json.dumps(serialisable, indent=2).encode("utf-8"))
st.sidebar.download_button("💾 Download config JSON", cfg_bytes.getvalue(),
                           file_name="qaqc_config.json", mime="application/json")
