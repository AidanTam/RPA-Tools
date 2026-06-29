# # blanks_app.py – Blanks QAQC Charts (dynamic, data-agnostic)
# import io, os, tempfile, datetime as dt
# import numpy as np
# import pandas as pd
# import matplotlib.pyplot as plt
# import matplotlib.dates as mdates
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches

# # ------------------------------------------------------------------
# # Page
# # ------------------------------------------------------------------
# st.set_page_config(page_title="Blanks QAQC Charts", layout="wide")
# st.title("🧼 Blanks QAQC Charts")

# # ------------------------------------------------------------------
# # 1. Upload
# # ------------------------------------------------------------------
# csv_file = st.sidebar.file_uploader("Upload blanks CSV", ["csv"])
# if not csv_file:
#     st.stop()
# df = pd.read_csv(csv_file)

# # ------------------------------------------------------------------
# # Helpers
# # ------------------------------------------------------------------
# def _norm(s: str) -> str:
#     return str(s).strip().lower().replace(" ", "").replace("-", "").replace("_", "")

# def guess_col(*names):
#     targets = {_norm(n) for n in names}
#     for c in df.columns:
#         if _norm(c) in targets:
#             return c
#     return None

# # ------------------------------------------------------------------
# # 2. Column mapping
# # ------------------------------------------------------------------
# with st.sidebar.expander("🗺️  Column mapping", expanded=True):
#     # required
#     lab_col = st.selectbox(
#         "Lab column",
#         df.columns,
#         index=int(df.columns.get_indexer([guess_col("Lab")])[0]) if guess_col("Lab") else 0,
#     )
#     elem_col = st.selectbox(
#         "Element column",
#         df.columns,
#         index=int(df.columns.get_indexer([guess_col("Element")])[0]) if guess_col("Element") else 0,
#     )
#     val_col = st.selectbox(
#         "Value column",
#         df.columns,
#         index=int(df.columns.get_indexer([guess_col("Value")])[0]) if guess_col("Value") else 0,
#     )

#     # optional
#     date_guess = guess_col("Date", "SampleDate", "AnalysisDate")
#     date_col = st.selectbox("Date column", ["None"] + list(df.columns),
#                             index=(1 + list(df.columns).index(date_guess)) if date_guess else 0)

#     unit_guess = guess_col("Unit", "Units")
#     unit_col = st.selectbox("Unit column", ["None"] + list(df.columns),
#                             index=(1 + list(df.columns).index(unit_guess)) if unit_guess else 0)

#     lod_guess = guess_col("LOD", "DL", "DetectionLimit", "MDL")
#     lod_col = st.selectbox("LOD / DL column (optional)", ["None"] + list(df.columns),
#                            index=(1 + list(df.columns).index(lod_guess)) if lod_guess else 0)

#     btype_guess = guess_col("Type", "BlankType", "QAQCType")
#     btype_col = st.selectbox("Blank-type column (optional)", ["None"] + list(df.columns),
#                              index=(1 + list(df.columns).index(btype_guess)) if btype_guess else 0)

# # ------------------------------------------------------------------
# # 3. Core filters
# # ------------------------------------------------------------------
# with st.sidebar.expander("🔎 Filter core columns"):
#     labs_sel = st.multiselect(
#         "Labs",
#         df[lab_col].dropna().unique(),
#         default=list(df[lab_col].dropna().unique()),
#     )
#     elems_sel = st.multiselect(
#         "Elements",
#         df[elem_col].dropna().unique(),
#         default=list(df[elem_col].dropna().unique()),
#     )
#     if btype_col != "None":
#         btypes_sel = st.multiselect(
#             "Blank types",
#             df[btype_col].dropna().unique(),
#             default=list(df[btype_col].dropna().unique()),
#         )
#     else:
#         btypes_sel = None

# # ------------------------------------------------------------------
# # 4. Extra categorical filters (x3)
# # ------------------------------------------------------------------
# cat_filters = []
# with st.sidebar.expander("🎛️  Extra categorical filters"):
#     for i in range(1, 4):
#         col = st.selectbox(f"Filter {i} – column", ["None"] + list(df.columns), key=f"cat_col_{i}")
#         if col != "None":
#             vals = st.multiselect("Values", df[col].dropna().unique(), key=f"cat_vals_{i}")
#             cat_filters.append((col, vals))

# # ------------------------------------------------------------------
# # 5. Date filter  (NaT-safe)
# # ------------------------------------------------------------------
# date_range = None
# if date_col != "None":
#     with st.sidebar.expander("🗓️  Date filter"):
#         df[date_col] = pd.to_datetime(df[date_col], errors="coerce")

#         date_series = df[date_col].dropna()
#         if date_series.empty:
#             st.info(f"No valid dates in '{date_col}' – date filter disabled.")
#         else:
#             dmin, dmax = date_series.min().date(), date_series.max().date()
#             date_range = st.date_input("Date range", (dmin, dmax))

# # ------------------------------------------------------------------
# # 6. Threshold (limit) definition
# # ------------------------------------------------------------------
# thresholds = {}
# with st.sidebar.expander("🎯 Blank limits"):
#     mode = st.radio("How to set limits?", ["Constant per element", "Factor × LOD column"])
#     if mode == "Constant per element":
#         for el in elems_sel:
#             # give a small default >0 to help spot mis-set zeros; tweak as you like
#             thresholds[el] = st.number_input(f"{el} limit", value=0.0, key=f"lim_{el}")
#     else:
#         factor = st.number_input("Factor × LOD", value=1.0)
#         if lod_col == "None":
#             st.warning("Select an LOD column (above) to use factor mode.")
#             st.stop()

# # ------------------------------------------------------------------
# # 7. Plot options
# # ------------------------------------------------------------------
# with st.sidebar.expander("📐 Plot options"):
#     log_y = st.checkbox("Log scale (Y)")
#     custom_ylim = st.checkbox("Custom Y-limits")
#     y_min = y_max = None
#     if custom_ylim:
#         y_min = st.number_input("Y min", value=0.0)
#         y_max = st.number_input("Y max", value=1.0)

# run = st.sidebar.button("🚀 Run analysis")
# if not run:
#     st.stop()

# # ------------------------------------------------------------------
# # 8. Outputs prep
# # ------------------------------------------------------------------
# ppt = Presentation()
# summary_rows = []

# def tmp_fig(fig):
#     t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(t.name, dpi=300, bbox_inches="tight")
#     plt.close(fig)
#     return t.name

# # ------------------------------------------------------------------
# # 9. Crunch
# # ------------------------------------------------------------------
# for el in elems_sel:
#     for lab in labs_sel:
#         sub = df[(df[lab_col] == lab) & (df[elem_col] == el)].copy()
#         if sub.empty:
#             continue

#         if btypes_sel is not None:
#             sub = sub[sub[btype_col].isin(btypes_sel)]

#         # Apply extra categorical filters
#         for col, vals in cat_filters:
#             if vals:
#                 sub = sub[sub[col].isin(vals)]

#         # Ensure numeric values
#         sub[val_col] = pd.to_numeric(sub[val_col], errors="coerce")

#         # Date filter
#         if date_col != "None":
#             sub = sub.dropna(subset=[date_col])
#             if date_range:
#                 start = pd.to_datetime(date_range[0])
#                 end = pd.to_datetime(date_range[1]) + pd.Timedelta(days=1)
#                 sub = sub[(sub[date_col] >= start) & (sub[date_col] < end)]

#         # Drop rows without values
#         sub = sub.dropna(subset=[val_col])
#         if sub.empty:
#             continue

#         # Determine limit (vector or scalar)
#         if mode == "Constant per element":
#             lim_scalar = float(thresholds.get(el, 0))
#             lim_series = pd.Series(lim_scalar, index=sub.index)
#         else:
#             lod_num = pd.to_numeric(sub[lod_col], errors="coerce") if lod_col != "None" else np.nan
#             lim_series = lod_num * factor
#             # fallback if all NaN
#             if lim_series.notna().any():
#                 lim_scalar = float(np.nanmedian(lim_series))
#             else:
#                 lim_scalar = 0.0
#                 lim_series = pd.Series(lim_scalar, index=sub.index)

#         # Failures (> limit rowwise)
#         failures = sub[sub[val_col] > lim_series]
#         total = len(sub)
#         nfail = len(failures)
#         fr = (nfail / total * 100) if total else 0.0

#         # ------------------------------------------------------------------
#         # Plot
#         # ------------------------------------------------------------------
#         fig, ax = plt.subplots(figsize=(10, 5))
#         # X
#         if date_col != "None":
#             x_all = sub[date_col]
#             x_fail = failures[date_col]
#         else:
#             x_all = range(total)
#             x_fail = failures.index

#         # Points
#         ax.scatter(x_all, sub[val_col], color="black", s=12, label="Values", linewidth=0.2)
#         ax.axhline(y=lim_scalar, linestyle="--", color="black", label="Limit", linewidth=1)
#         ax.scatter(x_fail, failures[val_col], color="red", s=14, label="Failures", alpha=0.7)

#         if date_col != "None":
#             ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
#             ax.xaxis.set_major_locator(mdates.AutoDateLocator())
#             plt.setp(ax.get_xticklabels(), rotation=45, ha="right")
#             ax.set_xlabel("Date")
#         else:
#             ax.set_xlabel("Sample #")

#         unit = ""
#         if unit_col != "None":
#             # first non-null
#             unit_series = sub[unit_col].dropna()
#             if not unit_series.empty:
#                 unit = unit_series.iloc[0]
#         ax.set_ylabel(f"{el} ({unit})" if unit else el)

#         ax.set_title(f"{lab} – {el} blanks")
#         ax.legend(loc="upper left", fontsize=8)
#         if log_y:
#             ax.set_yscale("log")
#         if custom_ylim:
#             ax.set_ylim(y_min, y_max)
#         ax.grid(False)

#         # Annotate
#         txt = f"Samples {total}\nFails {nfail}\nRate {fr:.1f}%"
#         ax.text(
#             0.98,
#             0.95,
#             txt,
#             transform=ax.transAxes,
#             va="top",
#             ha="right",
#             bbox=dict(facecolor="white", alpha=0.7, edgecolor="none"),
#             fontsize=8,
#         )

#         st.pyplot(fig)

#         # PPT
#         ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(
#             tmp_fig(fig), Inches(1), Inches(1), width=Inches(8), height=Inches(4.5)
#         )

#         # Summary row
#         summary_rows.append(
#             {
#                 "Lab": lab,
#                 "Element": el,
#                 "Samples": total,
#                 "Fails": nfail,
#                 "FailRate%": round(fr, 2),
#                 "LimitUsed": lim_scalar,
#             }
#         )

# # ------------------------------------------------------------------
# # 10. Summary + downloads
# # ------------------------------------------------------------------
# summary = pd.DataFrame(summary_rows)
# st.subheader("📊 Summary")
# st.dataframe(summary)

# xls = io.BytesIO()
# summary.to_excel(xls, index=False)

# ppt_bytes = io.BytesIO()
# ppt.save(ppt_bytes)

# st.download_button("Download summary.xlsx", xls.getvalue(), file_name="summary_blanks.xlsx")
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name="blank_plots.pptx")
# blanks_app.py – Blanks QAQC Charts


#### Version 2 #########
# blanks_app.py – Blanks QAQC Charts
import os, io, json, tempfile, datetime as dt, re
import numpy as np, pandas as pd
import matplotlib.pyplot as plt, matplotlib.dates as mdates
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(page_title="Blanks QAQC Charts", layout="wide")
st.title("🧼 Blanks QAQC Charts")

# ───────── helpers ─────────
def clean_numeric(s):
    if pd.api.types.is_numeric_dtype(s):
        return pd.to_numeric(s, errors="coerce")
    return pd.to_numeric(
        s.astype(str)
         .str.strip()
         .str.replace(r"[^\d\.\-eE+]", "", regex=True),
        errors="coerce",
    )

def tmp_fig(fig):
    t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(t.name, dpi=300, bbox_inches="tight"); plt.close(fig); return t.name

_safe_key_pat = re.compile(r"[^0-9a-zA-Z_]+")
def safe_key(*parts): return _safe_key_pat.sub("_", "_".join("" if p is None else str(p) for p in parts))

def idx_of(val, opts):
    """Return index of *val* in opts (list or pandas Index) else 0."""
    lst = list(opts)
    return lst.index(val) if val in lst else 0

def safe_defaults(defs, opts): return [d for d in (defs or []) if d in opts]

# ───────── config load ─────────
cfg_file = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
if cfg_file and not st.session_state.get("cfg_loaded"):
    st.session_state.update(json.load(cfg_file)); st.session_state.cfg_loaded=True; st.success("Config loaded!")

# ───────── CSV load ─────────
csv_bytes = st.sidebar.file_uploader("📂 Upload CSV", ["csv"], key="csv_bytes")
df=None
if csv_bytes is not None:
    df = pd.read_csv(csv_bytes); st.session_state["csv_path"]=""
elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
    df = pd.read_csv(st.session_state["csv_path"])
if df is None: st.info("Upload CSV or load config."); st.stop()

# ───────── column mapping ─────────
with st.sidebar.expander("🗺️ Column mapping", True):
    lab_col  = st.selectbox("Lab column", df.columns,  index=idx_of(st.session_state.get("lab_col",df.columns[0]), df.columns), key="lab_col")
    elem_col = st.selectbox("Element column", df.columns, index=idx_of(st.session_state.get("elem_col",df.columns[0]), df.columns), key="elem_col")
    val_col  = st.selectbox("Value column", df.columns,  index=idx_of(st.session_state.get("val_col",df.columns[0]), df.columns), key="val_col")
    date_col = st.selectbox("Date column", ["None"]+list(df.columns), index=idx_of(st.session_state.get("date_col","None"), ["None"]+list(df.columns)), key="date_col")
    unit_col = st.selectbox("Unit column", ["None"]+list(df.columns), index=idx_of(st.session_state.get("unit_col","None"), ["None"]+list(df.columns)), key="unit_col")
    lod_col  = st.selectbox("LOD/DL column", ["None"]+list(df.columns), index=idx_of(st.session_state.get("lod_col","None"), ["None"]+list(df.columns)), key="lod_col")
    type_col = st.selectbox("Blank‑type column", ["None"]+list(df.columns), index=idx_of(st.session_state.get("type_col","None"), ["None"]+list(df.columns)), key="type_col")

# ───────── filters & options (unchanged) ─────────
with st.sidebar.expander("🔎 Core filters"):
    labs_sel  = st.multiselect("Labs", df[lab_col].dropna().unique(), default=safe_defaults(st.session_state.get("labs_sel"), df[lab_col].unique()), key="labs_sel")
    elems_sel = st.multiselect("Elements", df[elem_col].dropna().unique(), default=safe_defaults(st.session_state.get("elems_sel"), df[elem_col].unique()), key="elems_sel")
    types_sel=None
    if type_col!="None":
        types_sel = st.multiselect("Blank types", df[type_col].dropna().unique(), default=safe_defaults(st.session_state.get("types_sel"), df[type_col].unique()), key="types_sel")

cat_filters,num_filters=[],[]
with st.sidebar.expander("🎛️ Categorical filters"):
    for i in range(1,4):
        col=st.selectbox(f"Cat {i}", ["None"]+list(df.columns), index=idx_of(st.session_state.get(f"cat_col_{i}","None"),["None"]+list(df.columns)), key=f"cat_col_{i}")
        if col!="None":
            vals=st.multiselect("Values", df[col].dropna().unique(), default=safe_defaults(st.session_state.get(f"cat_vals_{i}"),df[col].unique()), key=f"cat_vals_{i}")
            cat_filters.append((col, vals))

with st.sidebar.expander("🔢 Numerical filters"):
    num_cols=[c for c in df.columns if clean_numeric(df[c]).notna().any()]
    for i in range(1,3):
        col=st.selectbox(f"Num {i}", ["None"]+num_cols, index=idx_of(st.session_state.get(f"num_col_{i}","None"),["None"]+num_cols), key=f"num_col_{i}")
        if col!="None":
            s=clean_numeric(df[col]); lo,hi=float(s.min()), float(s.max())
            rng=st.slider("Range", lo, hi, value=st.session_state.get(f"num_rng_{i}",(lo,hi)), key=f"num_rng_{i}")
            num_filters.append((col, rng))

with st.sidebar.expander("🗓️ Date filter"):
    date_rng=None
    if date_col!="None":
        df[date_col]=pd.to_datetime(df[date_col], errors="coerce")
        dmin,dmax=df[date_col].min(), df[date_col].max()
        date_rng=st.date_input("Range", st.session_state.get("date_rng",(None if pd.isna(dmin) else dmin.date(), None if pd.isna(dmax) else dmax.date())), key="date_rng")

with st.sidebar.expander("🎯 Limits"):
    lim_mode = st.radio("Set limits by:", ["Constant per element","Factor × LOD"], key="lim_mode")
    const_limits={}
    if lim_mode=="Constant per element":
        for el in elems_sel:
            const_limits[el]=st.number_input(f"{el}", value=st.session_state.get(f"lim_{el}",0.0), key=f"lim_{el}")
    else:
        factor=st.number_input("Factor × LOD", value=float(st.session_state.get("lod_factor",1.0)), key="lod_factor")
        if lod_col=="None": st.warning("Select an LOD column for factor mode.")

with st.sidebar.expander("📐 Style"):
    point_size = st.number_input("Point size",1,value=int(st.session_state.get("p_size",8)),step=1,key="p_size")
    line_width = st.number_input("Line width",0.1,value=float(st.session_state.get("l_width",0.8)),step=0.1,format="%.2f",key="l_width")
    point_col  = st.color_picker("Point colour", st.session_state.get("p_col","#000000"), key="p_col")
    limit_col  = st.color_picker("Limit colour",  st.session_state.get("lim_col","#86A586"), key="lim_col")
    log_y      = st.checkbox("Log Y-axis", value=st.session_state.get("log_y",False), key="log_y")
    custom_ylim = st.checkbox("Custom Y-limits", value=st.session_state.get("custom_ylim",False), key="custom_ylim")
    if custom_ylim:
        y_axis_min = st.number_input("Y min", value=float(st.session_state.get("y_axis_min",0.0)), key="y_axis_min")
        y_axis_max = st.number_input("Y max", value=float(st.session_state.get("y_axis_max",1.0)), key="y_axis_max")
    else:
        y_axis_min = y_axis_max = None

# ───────── Run button ─────────
if st.sidebar.button("🚀 Run"): st.session_state.run_plots=True
if not st.session_state.get("run_plots"): st.info("Click **Run** to generate charts."); st.stop()

# ───────── crunch / plot (unchanged logic) ─────────
ppt, summary_rows = Presentation(), []

for el in elems_sel:
    for lab in labs_sel:
        base = df[(df[elem_col]==el)&(df[lab_col]==lab)]
        if base.empty: continue
        type_vals = types_sel if types_sel is not None else (base[type_col].dropna().unique() if type_col!="None" else [None])

        for typ in type_vals:
            sub = base if typ is None else base[base[type_col]==typ]
            for c,v in cat_filters: sub=sub[sub[c].isin(v)] if v else sub
            for c,(lo,hi) in num_filters: sub=sub[(clean_numeric(sub[c])>=lo)&(clean_numeric(sub[c])<=hi)]
            if date_col!="None" and date_rng:
                start,end=map(pd.to_datetime,date_rng); end+=pd.Timedelta(days=1)
                sub=sub[(sub[date_col]>=start)&(sub[date_col]<end)]

            sub[val_col]=clean_numeric(sub[val_col]); sub=sub.dropna(subset=[val_col])
            if sub.empty: continue

            if lim_mode=="Constant per element":
                limit=float(const_limits.get(el,0)); sub["Limit"]=limit
            else:
                sub["Limit"]=clean_numeric(sub[lod_col])*factor if lod_col!="None" else 0
                limit=float(np.nanmedian(sub["Limit"])) if sub["Limit"].notna().any() else 0

            fails=sub[sub[val_col]>sub["Limit"]]; total,nfail=len(sub),len(fails); rate=nfail/total*100
            container=st.container()
            with container:
                tick_key=safe_key("ticks", el, lab, typ)
                tick_count= st.slider("Date ticks",1,50, st.session_state.get(tick_key,10), key=tick_key) if date_col!="None" else None

                fig,ax=plt.subplots(figsize=(10,4))
                if date_col!="None":
                    x=sub[date_col]; min_dt,max_dt=x.min(),x.max()
                    tick_dates=pd.date_range(min_dt,max_dt,periods=tick_count) if min_dt!=max_dt else [min_dt]
                    ax.set_xticks(tick_dates); ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
                    plt.setp(ax.get_xticklabels(),rotation=45,ha="right")
                else: x=range(total); ax.set_xlabel("Sample #")

                ax.scatter(x,sub[val_col],s=point_size,color=point_col,label="Values")
                ax.axhline(limit,ls="--",lw=line_width,color=limit_col)
                if not fails.empty:
                    ax.scatter(fails[date_col] if date_col!="None" else fails.index,fails[val_col],s=point_size,color="red",label="Fails")
                if log_y: ax.set_yscale("log")
                if custom_ylim and y_axis_min is not None and y_axis_max is not None:
                    ax.set_ylim(y_axis_min, y_axis_max)
                unit=sub[unit_col].dropna().iloc[0] if unit_col!="None" and not sub[unit_col].dropna().empty else ""
                typ_lbl=typ if typ is not None else ""
                ax.set_ylabel(f"{el} ({unit})" if unit else el)
                ax.set_title(f"Blanks for {el} {unit} - {typ_lbl} - {lab}".strip())
                ax.legend(fontsize=8,loc="upper left")
                ax.text(0.99,0.95,f"Samples: {total}\nFails: {nfail}\nRate: {rate:.1f}%",
                        transform=ax.transAxes,ha="right",va="top",bbox=dict(fc="white",alpha=.8))
                ax.grid(False)
                st.pyplot(fig)
                ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(tmp_fig(fig),Inches(1),Inches(1),width=Inches(8),height=Inches(4.5))

            # summary_rows.append({"Lab":lab,"Element":el,"Type":typ_lbl,"Samples":total,"Fails":nfail,"Rate%":round(rate,2),"Limit":limit})
            # NEW block ─ copy‑paste this instead
            summary_rows.append({
                "Lab": lab,
                "Element": el,
                "Type": typ_lbl,
                "Samples": total,
                "Fails": nfail,
                "Rate%": round(rate, 2),
                "Limit": limit,
                "Min": round(sub[val_col].min(), 2),
                "Max": round(sub[val_col].max(), 2),
                "Mean": round(sub[val_col].mean(), 2),
            })

# ───────── summary/export ─────────
summary_df=pd.DataFrame(summary_rows)
st.subheader("📊 Summary"); st.dataframe(summary_df)
xls=io.BytesIO(); summary_df.to_excel(xls,index=False)
ppt_bytes=io.BytesIO(); ppt.save(ppt_bytes)
st.download_button("Download summary.xlsx",xls.getvalue(),file_name="summary_blanks.xlsx")
st.download_button("Download plots.pptx",ppt_bytes.getvalue(),file_name="blank_plots.pptx")

# ───────── save config ─────────
serialisable={}
for k,v in st.session_state.items():
    if k.startswith("_"): continue
    try: json.dumps(v); serialisable[k]=v
    except TypeError: pass
cfg_bytes=io.BytesIO(json.dumps(serialisable,indent=2).encode())
st.sidebar.download_button("💾 Download config JSON", cfg_bytes.getvalue(), file_name="qaqc_blanks_config.json", mime="application/json")
