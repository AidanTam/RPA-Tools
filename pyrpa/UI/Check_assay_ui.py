# ##### Streamlit app – Check Assay QAQC Charts
# """
# Author: Hamza Salhi – July 2025
# Convert of Jupyter notebook logic to a fully‑interactive Streamlit app, matching
# the style of the Blanks/CRM/Duplicates tools.
# """

# import os, io, json, tempfile, re, datetime as dt
# import numpy as np, pandas as pd
# import matplotlib.pyplot as plt, matplotlib.dates as mdates
# from matplotlib.ticker import FuncFormatter
# import streamlit as st
# from pptx import Presentation
# from pptx.util import Inches

# # ────────────────────────────────────────────────────────────
# # 🔧 Page & title
# # ────────────────────────────────────────────────────────────
# st.set_page_config(page_title="Check Assay QAQC Charts", layout="wide")
# st.title("🔍 Check Assay QAQC Charts")

# # ────────────────────────────────────────────────────────────
# # 🔧 Helpers (re‑used from other QAQC apps for consistency)
# # ────────────────────────────────────────────────────────────


# def clean_numeric(s):
#     """Try very hard to coerce column *s* (Series) to numeric."""
#     if pd.api.types.is_numeric_dtype(s):
#         return pd.to_numeric(s, errors="coerce")
#     return pd.to_numeric(
#         s.astype(str)
#         .str.strip()
#         .str.replace(r"[^\d\.\-eE+]", "", regex=True),
#         errors="coerce",
#     )


# def tmp_fig(fig):
#     """Save *fig* to a NamedTemporaryFile and return the path (PNG, 300 dpi)."""
#     t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
#     fig.savefig(t.name, dpi=300, bbox_inches="tight")
#     plt.close(fig)
#     return t.name


# _safe_key_pat = re.compile(r"[^0-9a-zA-Z_]+")

# def safe_key(*parts):
#     return _safe_key_pat.sub("_", "_".join("" if p is None else str(p) for p in parts))


# def idx_of(val, opts):
#     lst = list(opts)
#     return lst.index(val) if val in lst else 0


# def safe_defaults(defs, opts):
#     return [d for d in (defs or []) if d in opts]


# # ────────────────────────────────────────────────────────────
# # 🔧 Load / save config
# # ────────────────────────────────────────────────────────────
# cfg_file = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
# if cfg_file and not st.session_state.get("cfg_loaded"):
#     st.session_state.update(json.load(cfg_file))
#     st.session_state.cfg_loaded = True
#     st.success("Config loaded!")

# # ────────────────────────────────────────────────────────────
# # 📂 CSV upload / reload from session
# # ────────────────────────────────────────────────────────────
# csv_bytes = st.sidebar.file_uploader("📂 Upload CSV", ["csv"], key="csv_bytes")
# df = None
# if csv_bytes is not None:
#     df = pd.read_csv(csv_bytes)
#     st.session_state["csv_path"] = ""
# elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
#     df = pd.read_csv(st.session_state["csv_path"])
# if df is None:
#     st.info("Upload CSV or load config JSON to get started.")
#     st.stop()

# # ────────────────────────────────────────────────────────────
# # 🗺️ Column mapping
# # ────────────────────────────────────────────────────────────
# with st.sidebar.expander("🗺️ Column mapping", True):
#     elem_col = st.selectbox(
#         "Element column",
#         df.columns,
#         index=idx_of(st.session_state.get("elem_col", df.columns[0]), df.columns),
#         key="elem_col",
#     )
#     original_col = st.selectbox(
#         "Original value column",
#         df.columns,
#         index=idx_of(st.session_state.get("original_col", df.columns[1]), df.columns),
#         key="original_col",
#     )
#     duplicate_col = st.selectbox(
#         "Duplicate value column",
#         df.columns,
#         index=idx_of(st.session_state.get("duplicate_col", df.columns[2]), df.columns),
#         key="duplicate_col",
#     )
#     prim_lab_col = st.selectbox(
#         "Primary lab column",
#         df.columns,
#         index=idx_of(st.session_state.get("prim_lab_col", df.columns[3]), df.columns),
#         key="prim_lab_col",
#     )
#     sec_lab_col = st.selectbox(
#         "Secondary lab column",
#         df.columns,
#         index=idx_of(st.session_state.get("sec_lab_col", df.columns[4]), df.columns),
#         key="sec_lab_col",
#     )
#     dtype_col = st.selectbox(
#         "Duplicate‑type column",
#         ["None"] + list(df.columns),
#         index=idx_of(st.session_state.get("dtype_col", "None"), ["None"] + list(df.columns)),
#         key="dtype_col",
#     )
#     unit_col = st.selectbox(
#         "Unit column (optional)",
#         ["None"] + list(df.columns),
#         index=idx_of(st.session_state.get("unit_col", "None"), ["None"] + list(df.columns)),
#         key="unit_col",
#     )
#     date_col = st.selectbox(
#         "Date column (optional)",
#         ["None"] + list(df.columns),
#         index=idx_of(st.session_state.get("date_col", "None"), ["None"] + list(df.columns)),
#         key="date_col",
#     )

# # ────────────────────────────────────────────────────────────
# # 🔎 Core filters
# # ────────────────────────────────────────────────────────────
# with st.sidebar.expander("🔎 Core filters"):
#     elems_sel = st.multiselect(
#         "Elements",
#         df[elem_col].dropna().unique(),
#         default=safe_defaults(st.session_state.get("elems_sel"), df[elem_col].unique()),
#         key="elems_sel",
#     )
#     prim_labs_sel = st.multiselect(
#         "Primary labs",
#         df[prim_lab_col].dropna().unique(),
#         default=safe_defaults(st.session_state.get("prim_labs_sel"), df[prim_lab_col].unique()),
#         key="prim_labs_sel",
#     )
#     sec_labs_sel = st.multiselect(
#         "Secondary labs",
#         df[sec_lab_col].dropna().unique(),
#         default=safe_defaults(st.session_state.get("sec_labs_sel"), df[sec_lab_col].unique()),
#         key="sec_labs_sel",
#     )
#     dup_types_sel = None
#     if dtype_col != "None":
#         dup_types_sel = st.multiselect(
#             "Duplicate types",  # usually just "Check Assay"
#             df[dtype_col].dropna().unique(),
#             default=safe_defaults(st.session_state.get("dup_types_sel"), df[dtype_col].unique()),
#             key="dup_types_sel",
#         )

# # Additional categorical & numeric filters (copy pattern from Blanks app) ----------
# cat_filters, num_filters = [], []
# with st.sidebar.expander("🎛️ Categorical filters"):
#     for i in range(1, 4):
#         col = st.selectbox(
#             f"Cat {i}",
#             ["None"] + list(df.columns),
#             index=idx_of(st.session_state.get(f"cat_col_{i}", "None"), ["None"] + list(df.columns)),
#             key=f"cat_col_{i}",
#         )
#         if col != "None":
#             vals = st.multiselect(
#                 "Values",
#                 df[col].dropna().unique(),
#                 default=safe_defaults(
#                     st.session_state.get(f"cat_vals_{i}"), df[col].unique()
#                 ),
#                 key=f"cat_vals_{i}",
#             )
#             cat_filters.append((col, vals))

# with st.sidebar.expander("🔢 Numerical filters"):
#     num_cols = [c for c in df.columns if clean_numeric(df[c]).notna().any()]
#     for i in range(1, 3):
#         col = st.selectbox(
#             f"Num {i}",
#             ["None"] + num_cols,
#             index=idx_of(st.session_state.get(f"num_col_{i}", "None"), ["None"] + num_cols),
#             key=f"num_col_{i}",
#         )
#         if col != "None":
#             s = clean_numeric(df[col])
#             lo, hi = float(s.min()), float(s.max())
#             rng = st.slider(
#                 "Range",
#                 lo,
#                 hi,
#                 value=st.session_state.get(f"num_rng_{i}", (lo, hi)),
#                 key=f"num_rng_{i}",
#             )
#             num_filters.append((col, rng))

# with st.sidebar.expander("🗓️ Date filter"):
#     date_rng = None
#     if date_col != "None":
#         df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
#         dmin, dmax = df[date_col].min(), df[date_col].max()
#         date_rng = st.date_input(
#             "Range",
#             st.session_state.get(
#                 "date_rng",
#                 (
#                     None if pd.isna(dmin) else dmin.date(),
#                     None if pd.isna(dmax) else dmax.date(),
#                 ),
#             ),
#             key="date_rng",
#         )

# # Style options ------------------------------------------------------------------
# with st.sidebar.expander("📐 Style"):
#     point_size = st.number_input(
#         "Point size", 1, value=int(st.session_state.get("p_size", 5)), key="p_size"
#     )
#     line_width = st.number_input(
#         "Line width", 0.1, value=float(st.session_state.get("l_width", 0.8)), step=0.1, format="%.2f", key="l_width"
#     )
#     point_col = st.color_picker("Point colour", st.session_state.get("p_col", "#000000"), key="p_col")
#     one2one_col = st.color_picker("1:1 line colour", st.session_state.get("one1_col", "#808080"), key="one1_col")
#     log_axes = st.checkbox("Log scale", value=st.session_state.get("log_axes", True), key="log_axes")

# # ────────────────────────────────────────────────────────────
# # 🚀 Run button
# # ────────────────────────────────────────────────────────────
# if st.sidebar.button("🚀 Run"):
#     st.session_state.run_plots = True
# if not st.session_state.get("run_plots"):
#     st.info("Set your options and click **Run** to generate charts.")
#     st.stop()

# # ────────────────────────────────────────────────────────────
# # 🧮 Stats helper
# # ────────────────────────────────────────────────────────────

# def calc_stats(original: pd.Series, duplicate: pd.Series):
#     return {
#         "Count": len(original),
#         "Mean_or": original.mean(),
#         "Mean_dp": duplicate.mean(),
#         "Max_or": original.max(),
#         "Max_dp": duplicate.max(),
#         "Min_or": original.min(),
#         "Min_dp": duplicate.min(),
#         "Median_or": original.median(),
#         "Median_dp": duplicate.median(),
#         "Q3_or": original.quantile(0.75),
#         "Q3_dp": duplicate.quantile(0.75),
#         "Std_or": original.std(ddof=1),
#         "Std_dp": duplicate.std(ddof=1),
#         "CoefVar_or": original.std(ddof=1) / original.mean() * 100 if original.mean() else np.nan,
#         "CoefVar_dp": duplicate.std(ddof=1) / duplicate.mean() * 100 if duplicate.mean() else np.nan,
#         "Correlation": np.corrcoef(original, duplicate)[0, 1] if len(original) > 1 else np.nan,
#     }

# # ────────────────────────────────────────────────────────────
# # 🔄 Crunch & plot
# # ────────────────────────────────────────────────────────────
# ppt = Presentation()
# stats_rows = []

# # pre‑clean numeric value columns once
# for c in [original_col, duplicate_col]:
#     df[c] = clean_numeric(df[c])

# base_df = df.copy()

# # Apply core filters -------------------------------------------------------------
# if dup_types_sel is not None:
#     base_df = base_df[base_df[dtype_col].isin(dup_types_sel)]
# base_df = base_df[base_df[elem_col].isin(elems_sel)]
# base_df = base_df[base_df[prim_lab_col].isin(prim_labs_sel)]
# base_df = base_df[base_df[sec_lab_col].isin(sec_labs_sel)]
# for c, v in cat_filters:
#     if v:
#         base_df = base_df[base_df[c].isin(v)]
# for c, (lo, hi) in num_filters:
#     base_df = base_df[(clean_numeric(base_df[c]) >= lo) & (clean_numeric(base_df[c]) <= hi)]
# if date_col != "None" and date_rng:
#     start, end = map(pd.to_datetime, date_rng)
#     end += pd.Timedelta(days=1)
#     base_df = base_df[(base_df[date_col] >= start) & (base_df[date_col] < end)]

# if base_df.empty:
#     st.warning("No data after filtering – try relaxing filters.")
#     st.stop()

# # Iterate by element & lab pair --------------------------------------------------
# for el in elems_sel:
#     element_df = base_df[base_df[elem_col] == el]
#     if element_df.empty:
#         continue
#     pairs = element_df[[prim_lab_col, sec_lab_col]].drop_duplicates()

#     for _, pair in pairs.iterrows():
#         prim_lab = pair[prim_lab_col]
#         sec_lab = pair[sec_lab_col]
#         sub = element_df[
#             (element_df[prim_lab_col] == prim_lab) & (element_df[sec_lab_col] == sec_lab)
#         ]
#         if sub.empty:
#             continue

#         x = sub[original_col].dropna()
#         y = sub[duplicate_col].dropna()
#         common = sub.dropna(subset=[original_col, duplicate_col])
#         if common.empty:
#             continue

#         x = common[original_col]
#         y = common[duplicate_col]

#         # Percentiles for Q‑Q plot
#         percentiles = np.arange(1, 101)
#         xp = np.percentile(x, percentiles)
#         yp = np.percentile(y, percentiles)

#         # Statistics ------------------------------------------------------------
#         stats = calc_stats(x, y)
#         pretty_stats = {k: (round(v, 3) if isinstance(v, (int, float)) else v) for k, v in stats.items()}

#         # Plot ------------------------------------------------------------------
#         fig, axs = plt.subplots(1, 2, figsize=(14, 7))

#         # Q‑Q plot
#         axs[0].plot(xp, yp, "o", markersize=3, color=point_col, label="Percentiles")
#         axs[0].plot([xp.min(), xp.max()], [xp.min(), xp.max()], "--", lw=line_width, color=one2one_col)
#         axs[0].set_xlabel(f"{el} – Primary ({prim_lab})")
#         axs[0].set_ylabel(f"{el} – Secondary ({sec_lab})")
#         axs[0].set_title(f"Q‑Q plot – {el} ({prim_lab} vs {sec_lab})")

#         # Scatter plot
#         axs[1].scatter(x, y, s=point_size, color=point_col, edgecolors="black", linewidths=0.3, alpha=0.6)
#         axs[1].plot([x.min(), x.max()], [x.min(), x.max()], "--", lw=line_width, color=one2one_col)
#         axs[1].set_xlabel(f"{el} – Primary ({prim_lab})")
#         axs[1].set_ylabel(f"{el} – Secondary ({sec_lab})")
#         axs[1].set_title(f"Scatter – {el} ({prim_lab} vs {sec_lab})")

#         if log_axes:
#             for ax in axs:
#                 ax.set_xscale("log")
#                 ax.set_yscale("log")
#                 ax.xaxis.set_major_formatter(FuncFormatter(lambda v, _: f"{v:.1f}"))
#                 ax.yaxis.set_major_formatter(FuncFormatter(lambda v, _: f"{v:.1f}"))
#         for ax in axs:
#             ax.grid(False)

#         # stats box on first subplot
#         stats_text = "\n".join([f"{k}: {v}" for k, v in pretty_stats.items()])
#         props = dict(boxstyle="round", facecolor="lightgray", alpha=0.5)
#         axs[0].text(0.97, 0.03, stats_text, transform=axs[0].transAxes, fontsize=9, va="bottom", ha="right", bbox=props)

#         plt.tight_layout()
#         st.pyplot(fig)
#         # Save slide in PPT
#         ppt.slides.add_slide(ppt.slide_layouts[5]).shapes.add_picture(tmp_fig(fig), Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))

#         # summary row -----------------------------------------------------------
#         unit = (
#             common[unit_col].dropna().iloc[0] if unit_col != "None" and not common[unit_col].dropna().empty else ""
#         )
#         stats_rows.append(
#             {
#                 "Element": el,
#                 "Unit": unit,
#                 "Primary lab": prim_lab,
#                 "Secondary lab": sec_lab,
#                 **{k: (round(v, 3) if isinstance(v, (int, float)) else v) for k, v in stats.items()},
#             }
#         )

# # ────────────────────────────────────────────────────────────
# # 📊 Summary & exports
# # ────────────────────────────────────────────────────────────
# summary_df = pd.DataFrame(stats_rows)
# st.subheader("📊 Statistics summary")
# st.dataframe(summary_df)

# xls = io.BytesIO()
# summary_df.to_excel(xls, index=False)

# def_time = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
# file_base = f"check_assay_{def_time}"

# st.download_button("Download summary.xlsx", xls.getvalue(), file_name=f"{file_base}.xlsx")

# ppt_bytes = io.BytesIO()
# ppt.save(ppt_bytes)
# st.download_button("Download plots.pptx", ppt_bytes.getvalue(), file_name=f"{file_base}.pptx")

# # ────────────────────────────────────────────────────────────
# # 💾 Save config JSON
# # ────────────────────────────────────────────────────────────
# serialisable = {}
# for k, v in st.session_state.items():
#     if k.startswith("_"):
#         continue
#     try:
#         json.dumps(v)
#         serialisable[k] = v
#     except TypeError:
#         pass
# cfg_bytes = io.BytesIO(json.dumps(serialisable, indent=2).encode())
# st.sidebar.download_button("💾 Download config JSON", cfg_bytes.getvalue(), file_name="qaqc_check_assay_config.json", mime="application/json")


#### Version 2 ######

##### Streamlit app – Check‑Assay QAQC Charts
"""
Author: Hamza Salhi – July 2025
Interactive Streamlit version of the check‑assay QAQC notebook.

New in this build
-----------------
• Sidebar *Duplicate types* multiselect now controls:
  → filtering **and** plot splitting.
  Each selected type gets its own chart set.
  If nothing is selected, data aren’t split further.
• Type added to the summary table.
"""

# ───────── imports ─────────
import os, io, json, tempfile, re, datetime as dt
import numpy as np, pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pyrpa.UI import common

# ───────── page ─────────
st.set_page_config(page_title="Check‑Assay QAQC Charts", layout="wide")
st.title("🔍 Check‑Assay QAQC Charts")

# ───────── helpers ─────────
_safe_pat = re.compile(r"[^0-9A-Za-z_]+")
safe_key  = lambda *p: _safe_pat.sub("_", "_".join("" if x is None else str(x) for x in p))
idx_of    = lambda v, opts: list(opts).index(v) if v in opts else 0

def safe_defaults(defaults, options):
    return [d for d in (defaults or []) if d in options]

def clean_numeric(s):
    if pd.api.types.is_numeric_dtype(s):
        return pd.to_numeric(s, errors="coerce")
    return pd.to_numeric(
        s.astype(str)
         .str.strip()
         .str.replace(r"[^\d\.\-eE+]", "", regex=True),
        errors="coerce"
    )

def tmp_fig(fig):
    t = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(t.name, dpi=300, bbox_inches="tight"); plt.close(fig)
    return t.name

_fmt = FuncFormatter(lambda v, _: f"{v:.2f}")

# ───────── config load ─────────
cfg_up = st.sidebar.file_uploader("🔄 Load config (JSON)", ["json"])
if cfg_up and not st.session_state.get("cfg_loaded"):
    st.session_state.update(json.load(cfg_up))
    st.session_state.cfg_loaded = True
    st.success("Config loaded!")

# ───────── data load ─────────
csv_up = st.sidebar.file_uploader("📂 Upload data (CSV or Excel)", ["csv","xlsx","xls"], key="csv")
df = None
if csv_up is not None:
    df = common.read_data_file(csv_up)
    st.session_state["csv_path"] = ""
elif st.session_state.get("csv_path") and os.path.exists(st.session_state["csv_path"]):
    df = common.read_data_file(st.session_state["csv_path"])

if df is None:
    st.info("Upload a CSV/Excel file or load a config to begin.")
    st.stop()

# ───────── column mapping ─────────
with st.sidebar.expander("🗺️ Column mapping", True):
    prim_lab_col = st.selectbox("Primary‑lab column", df.columns,
                                index=idx_of(st.session_state.get("prim_lab_col", df.columns[0]), df.columns),
                                key="prim_lab_col")
    sec_lab_col  = st.selectbox("Secondary‑lab column", df.columns,
                                index=idx_of(st.session_state.get("sec_lab_col", df.columns[1]), df.columns),
                                key="sec_lab_col")
    elem_col     = st.selectbox("Element column", df.columns,
                                index=idx_of(st.session_state.get("elem_col", df.columns[2]), df.columns),
                                key="elem_col")
    or_col       = st.selectbox("Original‑assay column", df.columns,
                                index=idx_of(st.session_state.get("or_col", df.columns[3]), df.columns),
                                key="or_col")
    dup_col      = st.selectbox("Duplicate‑assay column", df.columns,
                                index=idx_of(st.session_state.get("dup_col", df.columns[4]), df.columns),
                                key="dup_col")
    unit_col     = st.selectbox("Unit column", ["None"] + list(df.columns),
                                index=idx_of(st.session_state.get("unit_col", "None"),
                                             ["None"] + list(df.columns)),
                                key="unit_col")
    type_col     = st.selectbox("Duplicate‑type column", ["None"] + list(df.columns),
                                index=idx_of(st.session_state.get("type_col", "None"),
                                             ["None"] + list(df.columns)),
                                key="type_col")
    date_col     = st.selectbox("Date column", ["None"] + list(df.columns),
                                index=idx_of(st.session_state.get("date_col", "None"),
                                             ["None"] + list(df.columns)),
                                key="date_col")

# ───────── filters ─────────
with st.sidebar.expander("🔎 Core filters"):
    elem_opts = df[elem_col].dropna().unique()
    elems_sel = st.multiselect("Elements", elem_opts,
                               default=safe_defaults(st.session_state.get("elems_sel"), elem_opts),
                               key="elems_sel")

    prim_opts = df[prim_lab_col].dropna().unique()
    prim_labs = st.multiselect("Primary labs", prim_opts,
                               default=safe_defaults(st.session_state.get("prim_labs"), prim_opts),
                               key="prim_labs")

    sec_opts = df[sec_lab_col].dropna().unique()
    sec_labs = st.multiselect("Secondary labs", sec_opts,
                              default=safe_defaults(st.session_state.get("sec_labs"), sec_opts),
                              key="sec_labs")

    dup_types_sel = None
    if type_col != "None":
        dup_opts = df[type_col].dropna().unique()
        dup_types_sel = st.multiselect("Duplicate types (split plots)", dup_opts,
                                       default=safe_defaults(st.session_state.get("dup_types"), dup_opts),
                                       key="dup_types")

with st.sidebar.expander("🗓️ Date filter"):
    date_rng = None
    if date_col != "None":
        df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
        dmin, dmax = df[date_col].min(), df[date_col].max()
        date_rng = st.date_input("Range",
                                 st.session_state.get("date_rng",
                                         (None if pd.isna(dmin) else dmin.date(),
                                          None if pd.isna(dmax) else dmax.date())),
                                 key="date_rng")

# ───────── style ─────────
with st.sidebar.expander("📐 Style"):
    point_size = st.number_input("Point size (pts)", 1,
                                 value=int(st.session_state.get("p_size", 6)), step=1,
                                 key="p_size")
    line_width = st.number_input("Line width", 0.1,
                                 value=float(st.session_state.get("l_width", 0.6)), step=0.1,
                                 format="%.2f", key="l_width")
    point_col  = st.color_picker("Point colour", st.session_state.get("p_col", "#000000"), key="p_col")
    diag_col   = st.color_picker("1:1 line colour", st.session_state.get("diag_col", "#86A586"), key="diag_col")
    log_axes   = st.checkbox("Log scales", value=st.session_state.get("log_axes", True), key="log_axes")

# ───────── run trigger ─────────
if st.sidebar.button("🚀 Run"): st.session_state.run_plots = True
if not st.session_state.get("run_plots"):
    st.info("Click **Run** to generate plots.")
    st.stop()

# ───────── crunch & plot ─────────
ppt = Presentation(); summary_rows = []
size_pts2 = point_size ** 2

for el in elems_sel:
    el_df = df[df[elem_col] == el]
    # figure out which duplicate-type buckets we’ll iterate over
    type_buckets = dup_types_sel if (type_col != "None" and dup_types_sel) else [None]

    for typ in type_buckets:
        typ_df = el_df if typ is None else el_df[el_df[type_col] == typ]
        if typ_df.empty: continue

        for plab in prim_labs:
            for slab in sec_labs:
                sub = typ_df[(typ_df[prim_lab_col] == plab) & (typ_df[sec_lab_col] == slab)]

                if date_col != "None" and date_rng:
                    start, end = map(pd.to_datetime, date_rng); end += pd.Timedelta(days=1)
                    sub = sub[(sub[date_col] >= start) & (sub[date_col] < end)]

                x = clean_numeric(sub[or_col]); y = clean_numeric(sub[dup_col])
                ok = ~x.isna() & ~y.isna(); x, y = x[ok], y[ok]
                if x.empty: continue

                unit = ""
                if unit_col != "None":
                    uvals = sub[unit_col].dropna()
                    unit  = uvals.iloc[0] if not uvals.empty else ""
                title_unit = f"{el} {unit}".strip()
                typ_lbl = typ if typ is not None else "All types"

                # Q-Q percentiles
                perc = np.arange(1, 101); xp = np.percentile(x, perc); yp = np.percentile(y, perc)

                fig, (ax0, ax1) = plt.subplots(1, 2, figsize=(13, 6))

                # Q‑Q
                ax0.scatter(xp, yp, s=size_pts2, c=point_col, edgecolors="none", alpha=0.7)
                ax0.plot([xp.min(), xp.max()], [xp.min(), xp.max()],
                         ls="--", lw=line_width, c=diag_col)
                if log_axes:
                    ax0.set_xscale("log"); ax0.set_yscale("log")
                    ax0.xaxis.set_major_formatter(_fmt); ax0.yaxis.set_major_formatter(_fmt)
                ax0.set_title(f"Q‑Q ({title_unit})\n{plab} vs {slab} – {typ_lbl}")
                ax0.set_xlabel(f"{title_unit} – {plab}")
                ax0.set_ylabel(f"{title_unit} – {slab}")

                # Scatter
                ax1.scatter(x, y, s=size_pts2, c=point_col, edgecolors="none", alpha=0.6)
                ax1.plot([x.min(), x.max()], [x.min(), x.max()],
                         ls="--", lw=line_width, c=diag_col)
                if log_axes:
                    ax1.set_xscale("log"); ax1.set_yscale("log")
                    ax1.xaxis.set_major_formatter(_fmt); ax1.yaxis.set_major_formatter(_fmt)
                ax1.set_title(f"Scatter ({title_unit})\n{plab} vs {slab} – {typ_lbl}")
                ax1.set_xlabel(f"{title_unit} – {plab}")
                ax1.set_ylabel(f"{title_unit} – {slab}")

                plt.tight_layout()

                # stats
                stats = {
                    "Count": len(x),
                    "Mean_or": np.mean(x), "Mean_dp": np.mean(y),
                    "Max_or": np.max(x),  "Max_dp": np.max(y),
                    "Min_or": np.min(x),  "Min_dp": np.min(y),
                    "Median_or": np.median(x), "Median_dp": np.median(y),
                    "Q3_or": np.percentile(x, 75), "Q3_dp": np.percentile(y, 75),
                    "Std_or": np.std(x, ddof=1),  "Std_dp": np.std(y, ddof=1),
                    "CoefVar_or": np.std(x, ddof=1) / np.mean(x) * 100,
                    "CoefVar_dp": np.std(y, ddof=1) / np.mean(y) * 100,
                    "Corr": np.corrcoef(x, y)[0, 1],
                }
                ax0.text(0.02, 0.98,
                         "\n".join(f"{k}: {v:.3f}" if k != "Count" else f"{k}: {v}"
                                   for k, v in stats.items()),
                         transform=ax0.transAxes, va="top", ha="left",
                         fontsize=8, bbox=dict(fc="white", alpha=.85, ec="none"))

                # Streamlit & PPT
                st.pyplot(fig, clear_figure=True)
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                slide.shapes.add_picture(tmp_fig(fig), Inches(0.5), Inches(0.8), width=Inches(9))

                summary_rows.append({
                    "Element": el, "Unit": unit,
                    "Type": typ_lbl,
                    "Primary Lab": plab, "Secondary Lab": slab,
                    **{k: round(v, 4) if k != "Count" else v for k, v in stats.items()}
                })

# ───────── summary / downloads ─────────
summary_df = pd.DataFrame(summary_rows)
st.subheader("📊 Statistics summary"); st.dataframe(summary_df)

xls = io.BytesIO(); summary_df.to_excel(xls, index=False)
st.download_button("Download summary.xlsx", xls.getvalue(),
                   file_name="check_assay_summary.xlsx")

ppt_bytes = io.BytesIO(); ppt.save(ppt_bytes)
st.download_button("Download plots.pptx", ppt_bytes.getvalue(),
                   file_name="check_assay_plots.pptx")

# ───────── save config ─────────
serial = {}
for k, v in st.session_state.items():
    if k.startswith("_"): continue
    try: json.dumps(v); serial[k] = v
    except TypeError: pass
cfg_bytes = io.BytesIO(json.dumps(serial, indent=2).encode())
st.sidebar.download_button("💾 Download config JSON",
                           cfg_bytes.getvalue(),
                           file_name="qaqc_checkassay_config.json",
                           mime="application/json")
